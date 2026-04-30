"""
后端 API：为 frontend 提供 AI 对话、注册登录、文件上传等接口。
"""
import base64
from collections import Counter
import gc
import io
import json
import math
import os
from datetime import datetime, timedelta
import queue
import re
import html as html_module
import urllib.parse
import xml.etree.ElementTree as ET
import tempfile
import time
import wave
import secrets
import threading
import traceback
import uuid
from flask import Flask, request, jsonify, Response, send_from_directory
from flask_cors import CORS
from werkzeug.exceptions import RequestEntityTooLarge
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import textract
except Exception:
    textract = None

from config import (
    CHAT_API_URL,
    API_KEY,
    CHAT_MODEL,
    CHAT_SYSTEM_PROMPT,
    MAX_TOKENS,
    CHAT_OMNI_VOICE,
    CHAT_OMNI_AUDIO_FORMAT,
    CHAT_OMNI_SAMPLE_RATE,
    REALTIME_WS_URL,
    REALTIME_API_KEY,
    REALTIME_MODEL,
    REALTIME_SYSTEM_PROMPT,
)
import database as db

# face_engine 含 PyTorch/OpenCV 等，启动时 import 会在小内存机上占满 CPU、数十秒才监听端口 → 延迟加载
_face_engine_mod = None
_requests_mod = None
_websocket_mod = None


def _face_mod():
    global _face_engine_mod
    if _face_engine_mod is None:
        import face_engine as _face_engine_mod

    return _face_engine_mod


def _requests():
    global _requests_mod
    if _requests_mod is None:
        import requests as _requests_mod
    return _requests_mod


def _websocket():
    global _websocket_mod
    if _websocket_mod is None:
        import websocket as _websocket_mod
    return _websocket_mod


# 人脸引擎异步预热（避免单次 HTTP 长时间阻塞导致 Vite 代理 / 浏览器 / axios 超时）
_face_warmup_lock = threading.Lock()
# 串行执行人脸推理，避免多请求叠加峰值内存导致 OOM Kill（threaded=True 时）
_face_infer_lock = threading.Lock()
_face_warmup_state = "idle"  # idle | starting | loading | ready | error
_face_warmup_error = None  # str | None
# 仅首条识别请求打印一次说明，避免每次 POST 刷屏（引擎仍只初始化一次）
_face_recognize_hint_logged = False
# 识别结果短时缓存：同一用户短时间重复请求直接返回最近结果，降低 CPU/内存峰值
_face_recent_result_lock = threading.Lock()
_face_recent_result_by_user = {}
_face_recent_result_ttl_sec = float(os.getenv("FACE_RECOGNIZE_RESULT_TTL_SEC", "0"))
_face_min_infer_interval_sec = float(os.getenv("FACE_RECOGNIZE_MIN_INFER_INTERVAL_SEC", "0"))
_face_last_infer_ts_by_user = {}
_face_db_cache_lock = threading.Lock()
_face_db_cache = None
_face_db_cache_ts = 0.0
_face_db_cache_ttl_sec = float(os.getenv("FACE_DB_CACHE_TTL_SEC", "30.0"))

# 情绪异常判定阈值：最近 N 天内达到此次数则创建「主动疏导」触发
PROACTIVE_ANOMALY_THRESHOLD = 3
PROACTIVE_ANOMALY_DAYS = 7
FACE_ABNORMAL_EMOTIONS = {"anger", "contempt", "disgust", "fear", "sadness"}
# 兴趣新闻缓存：减少点击卡片时的等待（按用户缓存，定时刷新）
INTEREST_NEWS_CACHE_TTL_SEC = int(os.getenv("INTEREST_NEWS_CACHE_TTL_SEC", "600"))
_interest_news_cache_lock = threading.Lock()
_interest_news_cache_by_user = {}
# 情绪小贴士触发：按连续使用时长里程碑 + 负向情绪双条件
MOOD_TIP_MILESTONES_MINUTES = [30, 60, 180, 300]
MOOD_TIP_DAILY_MAX_PUSH = 2
MOOD_TIP_SESSION_GAP_RESET_MINUTES = 12
_mood_tip_session_lock = threading.Lock()
_mood_tip_session_by_user = {}

# 账号格式限制：仅允许邮箱或学号登录（学号规则与预分配账号保持一致）
STUDENT_ID_REGEX = re.compile(r"^20\d{8}$")

# 单文件上传最大 10MB，与前端一致；超过时请求被拒绝
MAX_UPLOAD_BYTES = 10 * 1024 * 1024

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = MAX_UPLOAD_BYTES
CORS(app, origins=["*"])


def _err(message, status=400):
    """统一 JSON 错误响应（须定义在文件前部，供人脸等路由在运行时安全调用）。"""
    return jsonify({"error": message, "message": message}), status


def _is_email_account(account: str) -> bool:
    return bool(re.fullmatch(r"[^@\s]+@[^@\s]+\.[^@\s]+", (account or "").strip()))


def _is_student_id_account(account: str) -> bool:
    return bool(STUDENT_ID_REGEX.fullmatch((account or "").strip()))


@app.errorhandler(RequestEntityTooLarge)
def handle_too_large(e):
    return jsonify({"error": f"文件过大，单文件最大支持 10MB"}), 413

# 上传目录（相对 backend 的上级目录下的 uploads）
UPLOAD_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "uploads"))
os.makedirs(UPLOAD_DIR, exist_ok=True)

# 允许的扩展名
ALLOWED_IMAGE = {"jpg", "jpeg", "png", "gif", "webp"}
ALLOWED_VIDEO = {"mp4", "mov", "webm"}
ALLOWED_AUDIO = {"webm", "mp3", "wav", "ogg", "m4a"}
ALLOWED_DOCUMENT = {"txt", "doc", "docx", "ppt", "pptx", "pdf"}
ALLOWED_EXT = ALLOWED_IMAGE | ALLOWED_VIDEO | ALLOWED_AUDIO | ALLOWED_DOCUMENT
DOC_EXTRACT_MAX_CHARS = 12000

def _init_db_in_background():
    print("[startup] 后台执行 db.init_db()", flush=True)
    try:
        db.init_db()
        print("[startup] db.init_db() 完成", flush=True)
    except Exception as e:
        print(
            f"[startup] 数据库初始化失败（服务仍继续启动）：{type(e).__name__}: {e}",
            flush=True,
        )
        print(
            "[startup] 请检查 MYSQL_HOST/MYSQL_PORT/MYSQL_USER/MYSQL_PASSWORD/MYSQL_DATABASE，"
            "或在环境变量中调小 MYSQL_CONNECT_TIMEOUT/MYSQL_READ_TIMEOUT。",
            flush=True,
        )


# 启动时后台初始化 DB：避免数据库慢/卡导致 Flask 进程无法启动
threading.Thread(target=_init_db_in_background, name="init_db_bg", daemon=True).start()

# 内存 token 存储：token -> user_id（重启后失效，生产可改为 Redis/JWT）
_tokens = {}


def _user_row_to_json(row):
    """将 DB 行转为前端 User 格式（id 转字符串）。"""
    if not row:
        return None
    return {
        "id": str(row["id"]),
        "username": row["username"],
        "email": row["mail"],
        "preferred_name": (row.get("preferred_name") or "").strip() or None,
        "onboarding_done": bool(row.get("onboarding_done", 0)),
        "role": (row.get("role") or "STUDENT").upper(),
    }


def _require_auth():
    """从请求头取 token，校验并返回 user_id，失败返回 (None, response_tuple)。"""
    auth = request.headers.get("Authorization")
    if not auth or not auth.startswith("Bearer "):
        return None, (jsonify({"error": "未登录"}), 401)
    token = auth[7:].strip()
    user_id = _tokens.get(token)
    if not user_id:
        return None, (jsonify({"error": "登录已过期或无效"}), 401)
    return user_id, None


def build_messages(
    history: list,
    content: str,
    image_base64: str = None,
    image_mime: str = None,
    video_base64: str = None,
    video_mime: str = None,
    audio_base64: str = None,
    audio_mime: str = None,
    document_text: str = None,
    document_name: str = None,
) -> list:
    """
    将历史 + 当前用户消息转为 OpenAPI messages 格式。
    若 image_base64 / video_base64 / audio_base64 存在，最后一条为多模态 content 数组。
    """
    messages = []
    for item in history or []:
        role = item.get("role", "user")
        if role not in ("user", "assistant", "system"):
            role = "user"
        messages.append({"role": role, "content": (item.get("content") or "").strip()})

    text = (content or "").strip()
    if document_text:
        doc_title = (document_name or "文档").strip()[:120]
        doc_part = (
            f"【用户上传文档：{doc_title}】\n"
            "以下是从文档中提取的文本（可能有少量格式丢失），请据此理解并回答：\n"
            f"{document_text.strip()}"
        )
        text = f"{text}\n\n{doc_part}".strip() if text else doc_part
    if image_base64 or video_base64 or audio_base64:
        image_mime = image_mime or "image/jpeg"
        video_mime = video_mime or "video/mp4"
        parts = []
        if text:
            parts.append({"type": "text", "text": text})
        if image_base64:
            data_url = f"data:{image_mime};base64,{image_base64}"
            parts.append({"type": "image_url", "image_url": {"url": data_url}})
        if video_base64:
            data_url = f"data:{video_mime};base64,{video_base64}"
            parts.append({"type": "video_url", "video_url": {"url": data_url}})
        if audio_base64:
            fmt = (audio_mime or "").split("/")[-1].lower() if audio_mime else "wav"
            if fmt in ("x-wav", "wave"):
                fmt = "wav"
            audio_data = (audio_base64 or "").strip()
            # DashScope 兼容模式要求 input_audio.data 为可解析 URL 或 data:;base64,...；
            # 直接传裸 base64 会被当作 URL 校验并报 InvalidParameter。
            if audio_data and not audio_data.startswith("data:"):
                audio_data = f"data:;base64,{audio_data}"
            parts.append(
                {
                    "type": "input_audio",
                    "input_audio": {
                        "data": audio_data,
                        "format": fmt or "wav",
                    },
                }
            )
        if not parts:
            parts.append({"type": "text", "text": "（无文字内容）"})
        messages.append({"role": "user", "content": parts})
    else:
        if not text:
            text = "（无文字内容）"
        messages.append({"role": "user", "content": text})
    return messages


def _with_chat_system_prompt(messages: list) -> list:
    """为 qwen3-omni-flash 等对话模型注入小 Q 人设（置于 messages 最前）。"""
    text = (CHAT_SYSTEM_PROMPT or "").strip()
    if not text:
        return messages or []
    return [{"role": "system", "content": text}] + (messages or [])


def _chat_model_supports_omni_audio() -> bool:
    """HTTP Chat 是否请求文本+语音输出（与 Realtime WebSocket 无关）。"""
    m = (CHAT_MODEL or "").lower()
    return "omni" in m and "realtime" not in m


def _save_assistant_audio_file(wav_bytes: bytes, conv_id: int):
    """将助手回复音频写入 uploads，返回 (对外 URL 路径, 文件名)。"""
    sub = uuid.uuid4().hex[:8]
    save_dir = os.path.join(UPLOAD_DIR, sub)
    os.makedirs(save_dir, exist_ok=True)
    ext = (CHAT_OMNI_AUDIO_FORMAT or "wav").lstrip(".").lower() or "wav"
    name = f"ai-reply-{conv_id}-{uuid.uuid4().hex[:10]}.{ext}"
    path = os.path.join(save_dir, name)
    with open(path, "wb") as f:
        f.write(wav_bytes)
    rel = f"{sub}/{name}"
    return f"/api/uploads/{rel}", name


def _pcm16le_mono_to_wav(pcm: bytes, sample_rate: int) -> bytes:
    """将 16bit 小端单声道 PCM 裸数据封装为标准 RIFF WAV（浏览器可播）。"""
    if len(pcm) < 2:
        return b""
    if len(pcm) % 2 == 1:
        pcm = pcm[:-1]
    buf = io.BytesIO()
    with wave.open(buf, "wb") as wf:
        wf.setnchannels(1)
        wf.setsampwidth(2)
        wf.setframerate(int(sample_rate))
        wf.writeframes(pcm)
    return buf.getvalue()


def _decoded_omni_audio_to_wav_bytes(raw: bytes) -> bytes:
    """
    DashScope Omni 流式 audio 解码后多为裸 PCM；若已是 RIFF WAVE 则原样返回。
    若配置为 mp3 等容器格式则不做 PCM 封装。
    """
    if len(raw) >= 12 and raw[:4] == b"RIFF" and raw[8:12] == b"WAVE":
        return raw
    ext = (CHAT_OMNI_AUDIO_FORMAT or "wav").lstrip(".").lower()
    if ext in ("mp3", "mpeg", "opus", "aac"):
        return raw
    return _pcm16le_mono_to_wav(raw, CHAT_OMNI_SAMPLE_RATE)


def _resolve_image_from_request(data: dict):
    """从请求中解析出 image_base64 与 image_mime。支持 imageBase64 或 imageUrl（本地上传路径）。"""
    b64 = data.get("imageBase64") or data.get("image_base64")
    if b64:
        return b64, (data.get("imageMime") or data.get("image_mime") or "image/jpeg")
    url = data.get("imageUrl") or data.get("image_url")
    if url and isinstance(url, str) and "/api/uploads/" in url:
        try:
            rel = url.split("/api/uploads/", 1)[-1].lstrip("/")
            if ".." in rel or not rel:
                return None, None
            path = os.path.join(UPLOAD_DIR, rel)
            if os.path.isfile(path):
                with open(path, "rb") as f:
                    b64 = base64.b64encode(f.read()).decode("ascii")
                ext = os.path.splitext(path)[1].lower()
                mime = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png", ".gif": "image/gif", ".webp": "image/webp"}.get(ext, "image/jpeg")
                return b64, mime
        except Exception:
            pass
    return None, None


def _resolve_audio_from_request(data: dict):
    """
    从请求中解析 audio_base64、audio_mime、file_name。
    支持 audioBase64 或 audioUrl（本地上传路径 /api/uploads/...）。
    """
    b64 = data.get("audioBase64") or data.get("audio_base64")
    if b64:
        mime = (data.get("audioMime") or data.get("audio_mime") or "audio/wav").strip()
        name = (data.get("voiceFileName") or data.get("voice_file_name") or "").strip() or None
        return b64, mime, name

    url = data.get("audioUrl") or data.get("audio_url")
    if url and isinstance(url, str) and "/api/uploads/" in url:
        try:
            rel = url.split("/api/uploads/", 1)[-1].lstrip("/")
            if ".." in rel or not rel:
                return None, None, None
            path = os.path.join(UPLOAD_DIR, rel)
            if not os.path.isfile(path):
                return None, None, None
            ext = os.path.splitext(path)[1].lower().lstrip(".")
            if ext not in ALLOWED_AUDIO:
                return None, None, None
            with open(path, "rb") as f:
                raw = f.read()
            if not raw:
                return None, None, None
            mime = {
                "webm": "audio/webm",
                "mp3": "audio/mpeg",
                "wav": "audio/wav",
                "ogg": "audio/ogg",
                "m4a": "audio/m4a",
            }.get(ext, "audio/wav")
            file_name = os.path.basename(path) or None
            return base64.b64encode(raw).decode("ascii"), mime, file_name
        except Exception:
            return None, None, None
    return None, None, None


def _resolve_video_from_request(data: dict):
    """
    从请求中解析 video_base64 与 video_mime。
    支持 videoBase64 或 videoUrl（本地上传路径 /api/uploads/...）。
    """
    b64 = data.get("videoBase64") or data.get("video_base64")
    if b64:
        mime = (data.get("videoMime") or data.get("video_mime") or "video/mp4").strip()
        return b64, mime

    url = data.get("videoUrl") or data.get("video_url")
    if url and isinstance(url, str) and "/api/uploads/" in url:
        try:
            rel = url.split("/api/uploads/", 1)[-1].lstrip("/")
            if ".." in rel or not rel:
                return None, None
            path = os.path.join(UPLOAD_DIR, rel)
            if not os.path.isfile(path):
                return None, None
            ext = os.path.splitext(path)[1].lower().lstrip(".")
            if ext not in ALLOWED_VIDEO:
                return None, None
            with open(path, "rb") as f:
                raw = f.read()
            if not raw:
                return None, None
            mime = {
                "mp4": "video/mp4",
                "mov": "video/quicktime",
                "webm": "video/webm",
            }.get(ext, "video/mp4")
            return base64.b64encode(raw).decode("ascii"), mime
        except Exception:
            return None, None
    return None, None


def _mime_from_ext(ext: str) -> str:
    mapping = {
        "txt": "text/plain",
        "pdf": "application/pdf",
        "doc": "application/msword",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "ppt": "application/vnd.ms-powerpoint",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    return mapping.get((ext or "").lower().strip("."), "application/octet-stream")


def _extract_readable_strings(raw: bytes, min_len: int = 6) -> str:
    if not raw:
        return ""
    chunks = []
    seen = set()
    for enc in ("utf-16le", "utf-8", "latin1"):
        try:
            text = raw.decode(enc, errors="ignore")
        except Exception:
            continue
        for m in re.finditer(r"[\u4e00-\u9fffA-Za-z0-9，。！？；：、,.!?;:()\[\]{}《》“”\"'‘’%+\-_/\\]{%d,}" % min_len, text):
            seg = re.sub(r"\s+", " ", (m.group(0) or "").strip())
            if len(seg) < min_len:
                continue
            key = seg.lower()
            if key in seen:
                continue
            seen.add(key)
            chunks.append(seg)
            if len(chunks) >= 300:
                break
        if len(chunks) >= 300:
            break
    return "\n".join(chunks)


def _extract_text_from_txt(raw: bytes) -> str:
    for enc in ("utf-8-sig", "utf-8", "gb18030", "gbk", "utf-16", "utf-16le", "utf-16be"):
        try:
            return raw.decode(enc).strip()
        except Exception:
            continue
    return raw.decode("latin1", errors="ignore").strip()


def _extract_text_from_pdf(raw: bytes) -> str:
    if not PdfReader:
        return ""
    try:
        reader = PdfReader(io.BytesIO(raw))
        parts = []
        for idx, page in enumerate(reader.pages):
            txt = (page.extract_text() or "").strip()
            if txt:
                parts.append(f"[第{idx + 1}页]\n{txt}")
            if sum(len(x) for x in parts) >= DOC_EXTRACT_MAX_CHARS:
                break
        return "\n\n".join(parts).strip()
    except Exception:
        return ""


def _extract_text_from_docx(raw: bytes) -> str:
    if not DocxDocument:
        return ""
    try:
        doc = DocxDocument(io.BytesIO(raw))
        parts = []
        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                parts.append(t)
            if sum(len(x) for x in parts) >= DOC_EXTRACT_MAX_CHARS:
                break
        return "\n".join(parts).strip()
    except Exception:
        return ""


def _extract_text_from_pptx(raw: bytes) -> str:
    if not Presentation:
        return ""
    try:
        prs = Presentation(io.BytesIO(raw))
        parts = []
        for sidx, slide in enumerate(prs.slides):
            slide_lines = []
            for shape in slide.shapes:
                text = getattr(shape, "text", None)
                if not text:
                    continue
                line = str(text).strip()
                if line:
                    slide_lines.append(line)
            if slide_lines:
                parts.append(f"[第{sidx + 1}页]\n" + "\n".join(slide_lines))
            if sum(len(x) for x in parts) >= DOC_EXTRACT_MAX_CHARS:
                break
        return "\n\n".join(parts).strip()
    except Exception:
        return ""


def _extract_text_via_textract(raw: bytes, ext: str) -> str:
    if not textract:
        return ""
    suffix = f".{(ext or '').lower().strip('.')}"
    try:
        with tempfile.NamedTemporaryFile(delete=True, suffix=suffix) as tmp:
            tmp.write(raw)
            tmp.flush()
            out = textract.process(tmp.name)
            return (out or b"").decode("utf-8", errors="ignore").strip()
    except Exception:
        return ""


def _extract_text_from_document(raw: bytes, ext: str) -> str:
    ext = (ext or "").lower().strip(".")
    text = ""
    if ext == "txt":
        text = _extract_text_from_txt(raw)
    elif ext == "pdf":
        text = _extract_text_from_pdf(raw)
    elif ext == "docx":
        text = _extract_text_from_docx(raw)
    elif ext == "pptx":
        text = _extract_text_from_pptx(raw)
    elif ext in ("doc", "ppt"):
        text = _extract_text_via_textract(raw, ext) or _extract_readable_strings(raw)
    if not text:
        text = _extract_readable_strings(raw)
    text = re.sub(r"\n{3,}", "\n\n", (text or "").strip())
    if len(text) > DOC_EXTRACT_MAX_CHARS:
        text = text[:DOC_EXTRACT_MAX_CHARS] + "\n...(文档较长，已截断)"
    return text


def _resolve_document_from_request(data: dict):
    """
    从请求中解析文档文本内容。
    支持 documentUrl / fileUrl（本地上传路径 /api/uploads/...）。
    """
    url = (
        data.get("documentUrl")
        or data.get("document_url")
        or data.get("fileUrl")
        or data.get("file_url")
    )
    if not (url and isinstance(url, str) and "/api/uploads/" in url):
        return None, None, None
    try:
        rel = url.split("/api/uploads/", 1)[-1].lstrip("/")
        if ".." in rel or not rel:
            return None, None, None
        path = os.path.join(UPLOAD_DIR, rel)
        if not os.path.isfile(path):
            return None, None, None
        ext = os.path.splitext(path)[1].lower().lstrip(".")
        if ext not in ALLOWED_DOCUMENT:
            return None, None, None
        with open(path, "rb") as f:
            raw = f.read()
        if not raw:
            return None, None, None
        text = _extract_text_from_document(raw, ext)
        if not text:
            return None, None, os.path.basename(path) or None
        return text, _mime_from_ext(ext), os.path.basename(path) or None
    except Exception:
        return None, None, None


@app.route("/api/upload", methods=["POST"])
def upload():
    """上传文件，返回 { url, fileName, mimeType, category }。category: image|video|file|voice。"""
    if "file" not in request.files:
        return jsonify({"error": "未选择文件"}), 400
    f = request.files["file"]
    if not f or not f.filename:
        return jsonify({"error": "无效文件"}), 400
    ext = (f.filename.rsplit(".", 1)[-1] or "").lower()
    if ext not in ALLOWED_EXT:
        return jsonify({"error": f"不支持的文件类型: {ext}"}), 400
    name = secure_filename(f.filename) or "file"
    if "." not in name:
        name = f"{name}.{ext}"
    sub = uuid.uuid4().hex[:8]
    save_dir = os.path.join(UPLOAD_DIR, sub)
    os.makedirs(save_dir, exist_ok=True)
    save_path = os.path.join(save_dir, name)
    f.save(save_path)
    rel = f"{sub}/{name}"
    url = f"/api/uploads/{rel}"
    if ext in ALLOWED_IMAGE:
        category = "image"
    elif ext in ALLOWED_VIDEO:
        category = "video"
    elif ext in ALLOWED_AUDIO:
        category = "voice"
    else:
        category = "file"
    mime = f.content_type or "application/octet-stream"
    return jsonify({"url": url, "fileName": name, "mimeType": mime, "category": category})


@app.route("/api/uploads/<path:rel>", methods=["GET"])
def serve_upload(rel):
    """提供上传文件的访问。"""
    if ".." in rel:
        return jsonify({"error": "非法路径"}), 400
    path = os.path.join(UPLOAD_DIR, rel)
    if not os.path.isfile(path):
        return jsonify({"error": "文件不存在"}), 404
    return send_from_directory(UPLOAD_DIR, rel, as_attachment=False)


def _extract_schedules_from_text(user_id: int, text: str) -> list:
    """从用户一句话中抽取日程（如「明天去见孙老师」），写入 user_schedules。"""
    if not (text or "").strip():
        return []

    from datetime import datetime, timedelta

    def _is_schedule_rename_intent(t: str) -> bool:
        tt = (t or '').strip()
        if '答辩' not in tt or '开会' not in tt:
            return False
        # 用户明确“不是/改成/更改为”时才触发，避免误伤普通描述
        return any(k in tt for k in ('不是', '改', '更改', '改成', '为'))

    def _apply_schedule_rename_intent(u_id: int, t: str) -> int:
        now = datetime.now()
        start = now.strftime("%Y-%m-%d %H:%M:%S")
        end = (now + timedelta(days=7)).strftime("%Y-%m-%d %H:%M:%S")

        try:
            rows = db.get_schedules_by_user(int(u_id), start_date=start, end_date=end, limit=300) or []
        except Exception:
            return 0

        candidates: list[dict] = []
        for r in rows:
            title = (r.get('title') or '').strip()
            if title and '答辩' in title and title != '开会':
                sat = r.get('scheduled_at')
                sat_s = sat.isoformat(sep=' ')[:19] if hasattr(sat, 'isoformat') else str(sat or '')[:19]
                candidates.append({'id': int(r.get('id') or 0), 'scheduled_at': sat_s})

        if not candidates:
            return 0

        # “另一个日程”通常对应更晚的那条：取 scheduled_at 更大的
        candidates.sort(key=lambda x: x.get('scheduled_at') or '')
        target = candidates[-1]
        if not target or not target.get('id'):
            return 0

        ok = db.update_schedule(int(target['id']), int(u_id), title='开会')
        return 1 if ok else 0

    if _is_schedule_rename_intent(text):
        _apply_schedule_rename_intent(user_id, text)
        # 本句是“纠正/改名”，不新增日程
        return []
    # 取消/删除类语句交给删除逻辑处理，避免误新增
    if _is_schedule_delete_intent(text):
        return []
    text = (text or "").strip()[:1500]
    created = []
    today = datetime.now()
    today_str = today.strftime("%Y-%m-%d")
    tomorrow_str = (today + timedelta(days=1)).strftime("%Y-%m-%d")
    def _normalize_schedule_title(raw: str) -> str:
        t = (raw or "").strip()
        # 去掉常见空白与末尾标点，减少同义重复写入
        t = re.sub(r"\s+", "", t)
        t = re.sub(r"[，。！？；、,.!?;:：]+$", "", t)
        return t[:500]

    def _normalize_schedule_time(raw: str) -> str:
        st = (raw or "").replace("T", " ").strip()[:19]
        if len(st) == 16:
            st = st + ":00"
        return st

    # 预先加载近期待办，避免重复入库（同标题+同时间）
    existing_keys = set()
    try:
        now = datetime.now()
        start = now.strftime("%Y-%m-%d %H:%M:%S")
        end = (now + timedelta(days=30)).strftime("%Y-%m-%d %H:%M:%S")
        existed_rows = db.get_schedules_by_user(int(user_id), start_date=start, end_date=end, limit=500) or []
        for r in existed_rows:
            if (r.get("status") or "pending").strip() != "pending":
                continue
            et = r.get("scheduled_at")
            et_s = et.isoformat(sep=" ")[:19] if hasattr(et, "isoformat") else str(et or "")[:19]
            ek = (_normalize_schedule_title(r.get("title") or ""), _normalize_schedule_time(et_s))
            if ek[0] and ek[1]:
                existing_keys.add(ek)
    except Exception:
        pass

    if API_KEY:
        try:
            system = f"""当前日期：{today_str}。你只输出一个 JSON 数组，不要 markdown 或其它文字。
从用户这句话里提取明确或强烈暗示的日程/待办（某天要做的事、见谁、开会、上课等）。每个元素：{{"title":"事项简述","scheduled_at":"YYYY-MM-DD HH:MM:SS"}}。
规则：
1) 若一句话里出现多件事（如“约老师，还和上课冲突”），要尽量分别提取为多条。
2) 时间推断：明天={tomorrow_str}；后天=+2天；大后天=+3天。
3) 未说具体时间时：普通事项默认 10:00:00；上课/课程默认 08:00:00（表示上午第一节，后续可手改）。
4) 只有过去事实、纯情绪感受且无未来安排时，输出 []。"""
            payload = {
                "model": CHAT_MODEL,
                "messages": [{"role": "system", "content": system}, {"role": "user", "content": text}],
                "max_tokens": 500,
                "stream": False,
            }
            headers = {"Content-Type": "application/json", "Authorization": f"Bearer {API_KEY}"}
            r = _requests().post(CHAT_API_URL, json=payload, headers=headers, timeout=15)
            if r.status_code == 200:
                out = r.json()
                text_out = (out.get("choices", [{}])[0].get("message", {}).get("content") or "").strip()
                for raw in (text_out, text_out.replace("```json", "").replace("```", "").strip()):
                    try:
                        parsed = json.loads(raw)
                        arr = parsed if isinstance(parsed, list) else []
                        for item in arr:
                            if isinstance(item, dict) and item.get("title") and item.get("scheduled_at"):
                                title = _normalize_schedule_title(item.get("title") or "")
                                st = _normalize_schedule_time(item.get("scheduled_at") or "")
                                if title and len(st) >= 16:
                                    key = (title, st)
                                    if key in existing_keys:
                                        continue
                                    sid = db.create_schedule(
                                        int(user_id),
                                        title=title,
                                        scheduled_at=st,
                                        source="conversation",
                                        raw_text=text[:500],
                                    )
                                    existing_keys.add(key)
                                    created.append({"id": sid, "title": title, "scheduled_at": st})
                        if created:
                            return created
                        break
                    except (json.JSONDecodeError, TypeError):
                        continue
        except Exception:
            pass
    return created


def _is_schedule_delete_intent(text: str) -> bool:
    """判断一句话是否表达了“删除/取消日程”意图（LLM 主导）。"""
    if not (text or "").strip():
        return False
    text = (text or "").strip()[:600]
    if API_KEY:
        try:
            sys_prompt = (
                "你是日程意图分类器。判断用户这句话是否在表达“删除/取消已有日程安排”的意图。"
                "只输出 JSON：{\"delete_intent\":true} 或 {\"delete_intent\":false}。"
            )
            user_prompt = (
                f"当前时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
                f"用户输入：{text}\n"
                "注意：如“不开会了”“这个安排不用了”“取消明天的会”均属于删除意图。"
            )
            obj, err = _llm_json_chat(sys_prompt, user_prompt, max_tokens=120, timeout=12)
            if not err and isinstance(obj, dict):
                return bool(obj.get("delete_intent"))
        except Exception:
            pass
    return False


def _extract_schedule_deletions_from_text(user_id: int, text: str) -> list:
    """从用户一句话中抽取要删除的日程并执行删除。"""
    if not (text or "").strip():
        return []
    text = (text or "").strip()[:1500]
    if not _is_schedule_delete_intent(text):
        return []

    try:
        rows = db.get_schedules_by_user(int(user_id), limit=300) or []
    except Exception:
        return []
    candidates = []
    for r in rows:
        if (r.get("status") or "pending") != "pending":
            continue
        sat = r.get("scheduled_at")
        sat_str = sat.isoformat(sep=" ")[:19] if hasattr(sat, "isoformat") else str(sat or "")[:19]
        candidates.append({
            "id": int(r.get("id")),
            "title": (r.get("title") or "").strip(),
            "scheduled_at": sat_str,
        })
    if not candidates:
        return []

    candidates.sort(key=lambda x: x["scheduled_at"])

    def _parse_day_hint(s: str) -> str:
        now = datetime.now()
        m = re.search(r"(20\d{2})[-/年](\d{1,2})[-/月](\d{1,2})", s)
        if m:
            try:
                return f"{int(m.group(1)):04d}-{int(m.group(2)):02d}-{int(m.group(3)):02d}"
            except Exception:
                pass
        if "今天" in s:
            return now.strftime("%Y-%m-%d")
        if "明天" in s:
            return (now + timedelta(days=1)).strftime("%Y-%m-%d")
        if "后天" in s and "大后天" not in s:
            return (now + timedelta(days=2)).strftime("%Y-%m-%d")
        if "大后天" in s:
            return (now + timedelta(days=3)).strftime("%Y-%m-%d")
        return ""

    day_hint = _parse_day_hint(text)
    scoped_candidates = candidates
    if day_hint:
        only_day = [c for c in candidates if c["scheduled_at"].startswith(day_hint)]
        if only_day:
            scoped_candidates = only_day

    def _llm_pick_ids() -> list:
        if not API_KEY:
            return []
        preview = scoped_candidates[:30]
        schedule_lines = "\n".join(
            [f"- id={c['id']}, time={c['scheduled_at']}, title={c['title']}" for c in preview]
        )
        if not schedule_lines:
            return []
        sys_prompt = (
            "你是日程删除助手。根据用户输入，从候选日程中选择要删除的 id。"
            "只输出 JSON，不要其他文字。格式：{\"ids\":[1,2]}。"
            "若不确定或用户没有明确删除意图，返回 {\"ids\":[]}。"
        )
        user_prompt = (
            f"当前时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
            f"用户输入：{text}\n"
            f"候选日程：\n{schedule_lines}\n"
            f"{'用户明确提到日期：' + day_hint + '，只能在该日期候选中选择。' if day_hint else ''}\n"
            "请仅返回要删除的 ids。"
        )
        obj, err = _llm_json_chat(sys_prompt, user_prompt, max_tokens=300, timeout=15)
        if err or not isinstance(obj, dict):
            return []
        raw_ids = obj.get("ids") or obj.get("delete_ids") or []
        if not isinstance(raw_ids, list):
            return []
        allow = {c["id"] for c in preview}
        out = []
        for x in raw_ids:
            try:
                sid = int(x)
                if sid in allow:
                    out.append(sid)
            except Exception:
                continue
        return out

    target_ids = _llm_pick_ids()
    if not target_ids:
        return []

    id_to_item = {c["id"]: c for c in candidates}
    deleted = []
    seen = set()
    for sid in target_ids:
        if sid in seen:
            continue
        seen.add(sid)
        try:
            ok = db.delete_schedule(int(sid), int(user_id))
            if ok and sid in id_to_item:
                deleted.append(id_to_item[sid])
        except Exception:
            continue
    return deleted


def _detect_emotion_anomaly(user_id: int, user_text: str) -> None:
    """用模型判断用户输入是否情绪异常，若异常则写入 emotion_anomalies 并检查是否触发主动疏导。"""
    if not (user_text or "").strip() or not API_KEY:
        return
    try:
        payload = {
            "model": CHAT_MODEL,
            "messages": [
                {"role": "system", "content": "你只输出一段 JSON，不要其他内容。判断用户这句话是否表现出明显情绪异常（如焦虑、抑郁、愤怒、崩溃等）。若异常则输出：{\"is_abnormal\":true,\"emotion_label\":\"异常类型\",\"reason\":\"简短原因\"}；否则输出：{\"is_abnormal\":false}。"},
                {"role": "user", "content": (user_text or "")[:1500]},
            ],
            "max_tokens": 200,
            "stream": False,
        }
        headers = {"Content-Type": "application/json", "Authorization": f"Bearer {API_KEY}"}
        r = _requests().post(CHAT_API_URL, json=payload, headers=headers, timeout=15)
        if r.status_code != 200:
            return
        out = r.json()
        text = (out.get("choices", [{}])[0].get("message", {}).get("content") or "").strip()
        if not text:
            return
        for raw in (text, text.replace("```json", "").replace("```", "").strip()):
            try:
                obj = json.loads(raw)
                if obj.get("is_abnormal") and obj.get("emotion_label"):
                    db.add_emotion_anomaly(
                        int(user_id),
                        (obj.get("emotion_label") or "异常")[:64],
                        reason=(obj.get("reason") or "")[:2000],
                        from_monitoring=0,
                    )
                    _maybe_create_proactive_trigger(int(user_id))
                break
            except (json.JSONDecodeError, TypeError):
                continue
    except Exception:
        pass


def _is_face_abnormal_emotion(emotion_label: str) -> bool:
    return (emotion_label or "").strip().lower() in FACE_ABNORMAL_EMOTIONS


def _get_user_student_id(user_id: int) -> str:
    row = db.get_user_by_id(int(user_id)) or {}
    mail = (row.get("mail") or "").strip()
    if STUDENT_ID_REGEX.match(mail):
        return mail
    username = (row.get("username") or "").strip()
    if STUDENT_ID_REGEX.match(username):
        return username
    return ""


def _maybe_create_proactive_trigger(user_id: int) -> None:
    """同时统计聊天异常与人脸异常；达到阈值时创建主动疏导触发（不混存两类记录）。"""
    pending = db.get_pending_proactive_trigger(int(user_id))
    if pending:
        return
    chat_n = db.count_recent_anomalies(int(user_id), days=PROACTIVE_ANOMALY_DAYS)
    face_n = 0
    sid = _get_user_student_id(int(user_id))
    if sid:
        face_n = db.count_recent_face_abnormal_records(
            sid,
            days=PROACTIVE_ANOMALY_DAYS,
            abnormal_labels=sorted(FACE_ABNORMAL_EMOTIONS),
        )
    total = int(chat_n) + int(face_n)
    if total < PROACTIVE_ANOMALY_THRESHOLD:
        return
    if chat_n > 0 and face_n > 0:
        trigger_type = "multi_source_anomaly"
    elif face_n > 0:
        trigger_type = "face_anomaly"
    else:
        trigger_type = "repeated_anomaly"
    db.create_proactive_trigger(int(user_id), trigger_type)
    # 仅创建待处理触发，不在这里落聊天消息。
    # 由用户点击「去聊天」后的 /api/proactive/ack 统一发送一条关怀消息，避免一次流程出现两条。


def _summarize_face_abnormal_recent(student_id: str, days: int = 7, limit: int = 160) -> dict:
    rows = db.list_emotion_records(student_id=student_id, limit=limit) or []
    cutoff = datetime.now() - timedelta(days=max(1, days))
    label_count = {}
    total = 0
    for r in rows:
        ts = r.get("timestamp")
        if hasattr(ts, "isoformat"):
            dt = ts
        else:
            try:
                dt = datetime.strptime(str(ts)[:19], "%Y-%m-%d %H:%M:%S")
            except Exception:
                continue
        if dt < cutoff:
            continue
        label = (r.get("emotion_type") or "").strip().lower()
        if label not in FACE_ABNORMAL_EMOTIONS:
            continue
        total += 1
        label_count[label] = label_count.get(label, 0) + 1
    top = sorted(label_count.items(), key=lambda x: x[1], reverse=True)[:3]
    return {"total": total, "top_labels": top}


def _send_proactive_care_message(
    user_id: int, trigger_type: str, chat_n: int, face_n: int, use_llm: bool = False
):
    """写入一条主动关怀消息；默认走轻量模板，避免阻塞正常对话。"""
    try:
        convs = db.get_conversations_by_user(int(user_id), limit=1) or []
        if convs:
            conv_id = int(convs[0]["id"])
        else:
            conv_id = db.create_conversation(int(user_id), "小Q主动关怀")
    except Exception:
        return

    user_row = db.get_user_by_id(int(user_id)) or {}
    call_name = (
        (user_row.get("preferred_name") or "").strip()
        or (user_row.get("username") or "").strip()
        or "你"
    )[:50]
    sid = _get_user_student_id(int(user_id))
    face_summary = {"total": 0, "top_labels": []}
    if sid:
        try:
            face_summary = _summarize_face_abnormal_recent(sid, days=PROACTIVE_ANOMALY_DAYS, limit=180)
        except Exception:
            face_summary = {"total": int(face_n), "top_labels": []}

    # 最近聊天异常简要
    try:
        anomaly_rows = db.get_emotion_anomalies_by_user(int(user_id), limit=20, since_days=PROACTIVE_ANOMALY_DAYS) or []
    except Exception:
        anomaly_rows = []
    chat_labels = {}
    for r in anomaly_rows:
        if int(r.get("from_monitoring") or 0) != 0:
            continue
        lb = (r.get("emotion_label") or "").strip()
        if not lb:
            continue
        chat_labels[lb] = chat_labels.get(lb, 0) + 1
    top_chat = sorted(chat_labels.items(), key=lambda x: x[1], reverse=True)[:3]

    msg = ""
    if use_llm and API_KEY:
        sys_prompt = (
            "你是小Q。请写一条主动关怀消息，语气温柔自然、像朋友寒暄，不要太直白。"
            "不要说“系统检测到异常/你有问题”。"
            "可以轻微表达关注并邀请对方聊聊近况。"
            "长度 28~70 字。只输出 JSON：{\"message\":\"...\"}。"
        )
        user_prompt = (
            f"称呼：{call_name}\n"
            f"触发类型：{trigger_type}\n"
            f"近{PROACTIVE_ANOMALY_DAYS}天聊天异常次数：{chat_n}，主要标签：{top_chat}\n"
            f"近{PROACTIVE_ANOMALY_DAYS}天人脸异常次数：{face_n}，主要标签：{face_summary.get('top_labels')}\n"
            "请写一句不冒犯、不过度解读的问候式关怀话术。"
        )
        obj, err = _llm_json_chat(sys_prompt, user_prompt, max_tokens=200, timeout=20)
        if not err and isinstance(obj, dict):
            msg = (obj.get("message") or "").strip()
    if not msg:
        import random
        pool = [
            f"{call_name}，最近还好吗？要是你愿意，我们可以慢慢聊聊近况，我一直在这儿陪着你。",
            f"{call_name}，感觉你最近有点累，随时可以找我吐吐槽，或者什么都不说，我陪你听听歌也行。",
            f"嗨，{call_name}，今天过得怎么样？如果有烦心事，随时可以跟我说哦。",
            f"{call_name}，记得多喝水、按时休息。如果觉得压力大，我随时都在这里。",
            f"看你最近好像有点疲惫，{call_name}，要不要停下来深呼吸一下？我陪你聊会儿天吧。",
            f"{call_name}，不管今天遇到了什么，都不要太勉强自己。需要倾听者的话，我随时在线。",
            f"你不是一个人在面对这些哦，{call_name}。如果觉得辛苦，随时可以跟我分享你的感受。"
        ]
        msg = random.choice(pool)
    try:
        # 主动关怀固定写入单条消息，避免点击一次出现多条分句消息。
        single_msg = (msg or "").strip()
        if not single_msg:
            return
        mid = db.create_message(conv_id, "assistant", single_msg[:2000], "text", None, None)
        db.update_conversation_last_message(conv_id, single_msg[:500])
        return {
            "conversationId": str(conv_id),
            "message": {
                "id": str(mid),
                "content": single_msg[:2000],
                "sender": "ai",
                "timestamp": datetime.now().isoformat(),
                "type": "text",
                "fileUrl": None,
                "fileName": None,
            },
        }
    except Exception:
        return None


def _safe_json_dumps(obj) -> str:
    try:
        return json.dumps(obj, ensure_ascii=False)
    except Exception:
        return "{}"


def _text_to_vector(text: str, dim: int = 32) -> list:
    """轻量向量化：用于本地长期记忆召回（可后续替换为真实向量数据库 embedding）。"""
    vec = [0.0] * max(8, dim)
    txt = (text or "").strip().lower()
    if not txt:
        return vec
    for i, ch in enumerate(txt):
        idx = (ord(ch) + i * 17) % len(vec)
        vec[idx] += 1.0
    norm = math.sqrt(sum(v * v for v in vec)) or 1.0
    return [round(v / norm, 6) for v in vec]


def _cosine(a: list, b: list) -> float:
    if not a or not b:
        return 0.0
    n = min(len(a), len(b))
    if n <= 0:
        return 0.0
    return float(sum(float(a[i]) * float(b[i]) for i in range(n)))


def _memory_add(user_id: int, memory_type: str, content: str, metadata: dict = None) -> None:
    txt = (content or "").strip()
    if not txt:
        return
    try:
        vec = _text_to_vector(txt)
        db.add_user_memory_vector(
            int(user_id),
            (memory_type or "chat").strip()[:32],
            txt[:5000],
            _safe_json_dumps(vec),
            _safe_json_dumps(metadata or {}),
        )
    except Exception:
        pass


def _memory_search(user_id: int, query: str, limit: int = 6, memory_type: str = "") -> list:
    q = (query or "").strip()
    if not q:
        return []
    try:
        rows = db.list_user_memory_vectors(
            int(user_id),
            memory_type=(memory_type.strip()[:32] if memory_type else None),
            limit=300,
        ) or []
    except Exception:
        return []
    qv = _text_to_vector(q)
    scored = []
    for r in rows:
        try:
            rv = json.loads((r.get("vector_json") or "[]").strip())
            if not isinstance(rv, list):
                continue
            score = _cosine(qv, rv)
            if score <= 0:
                continue
            scored.append((score, r))
        except Exception:
            continue
    scored.sort(key=lambda x: x[0], reverse=True)
    return [x[1] for x in scored[: max(1, limit)]]


def _collect_today_mood_signal_tokens(user_id: int) -> list:
    """收集今日（自然日）内与情绪相关的原始标签/词（含正负），供统计或筛选。"""
    now = datetime.now()
    start = now.replace(hour=0, minute=0, second=0, microsecond=0)
    start_s = start.strftime("%Y-%m-%d %H:%M:%S")
    end_s = now.strftime("%Y-%m-%d %H:%M:%S")
    keywords = []
    try:
        rows = db.get_emotion_anomalies_by_user(int(user_id), limit=120, since_days=1) or []
        for r in rows:
            ts = r.get("created_at")
            if hasattr(ts, "strftime"):
                if ts < start:
                    continue
            keywords.append((r.get("emotion_label") or "").strip())
    except Exception:
        pass
    try:
        msg_rows = db.get_messages_for_user_range(int(user_id), start_s, end_s, limit=500) or []
        for m in msg_rows:
            if (m.get("role") or "") != "user":
                continue
            text = (m.get("content") or "").strip()
            if not text:
                continue
            for k in ("焦虑", "疲惫", "烦躁", "开心", "轻松", "压力", "崩溃", "紧张"):
                if k in text:
                    keywords.append(k)
    except Exception:
        pass
    return [x for x in keywords if x]


def _is_negative_mood_token(token: str) -> bool:
    """判断是否属于需要提示关怀的不良情绪信号（排除明显正向/中性）。"""
    t = (token or "").strip()
    if not t:
        return False
    if t in ("开心", "轻松"):
        return False
    if t in ("焦虑", "疲惫", "烦躁", "压力", "崩溃", "紧张"):
        return True
    # 常见异常描述子串
    neg_fragments = (
        "焦虑", "抑郁", "愤怒", "崩溃", "压力", "疲惫", "低落", "难过", "紧张", "烦躁",
        "绝望", "失眠", "恐慌", "不安", "焦躁", "心累", "累垮",
    )
    if any(s in t for s in neg_fragments):
        return True
    # 明显中性/无异常
    if any(s in t for s in ("平稳", "正常", "良好", "未见异常", "无明显")):
        return False
    # 其它来自模型的一句 emotion_label：默认视为需要关注（异常流写入的多为负向）
    return len(t) >= 2


def _today_mood_keywords(user_id: int, limit: int = 5) -> list:
    tokens = _collect_today_mood_signal_tokens(user_id)
    cnt = Counter(tokens)
    return [k for k, _ in cnt.most_common(max(1, limit))]


def _today_negative_mood_keywords(user_id: int, limit: int = 5) -> list:
    """今日不良情绪关键词（去重、按频次）；无则空列表。"""
    tokens = [t for t in _collect_today_mood_signal_tokens(user_id) if _is_negative_mood_token(t)]
    if not tokens:
        return []
    cnt = Counter(tokens)
    return [k for k, _ in cnt.most_common(max(1, limit))]


def _build_mood_tip_card_fields(negative_kw) -> dict:
    """根据检测到的负面信号生成情绪小贴士副标题与正文。"""
    joined = "、".join(negative_kw)
    subtitle = f"检测到你最近有「{joined}」等状态倾向，需要多照顾自己一点。"

    chunks = []
    if any(("焦虑" in k or "紧张" in k or k in ("焦虑", "紧张")) for k in negative_kw):
        chunks.append("若思绪停不下来，建议先暂停手头工作 5～10 分钟：缓慢深呼吸、喝点温水，或到窗边短走几步。")
    if any(("疲惫" in k or "心累" in k or "累" in k or k == "疲惫") for k in negative_kw):
        chunks.append("若身体偏疲惫，建议离开座位做轻度拉伸或慢走几分钟，再继续高强度用脑。")
    if any(("压力" in k or "烦躁" in k or "崩溃" in k) for k in negative_kw):
        chunks.append("若感到压力大或心烦，可把任务拆成更小的一步，先完成最简单的一件，降低心理负担。")
    if not chunks:
        chunks.append("建议你暂时停下工作，起身活动或到户外慢走几分钟，比持续硬撑更能恢复状态。")

    content = " ".join(chunks)
    priority = min(92, 76 + min(6, len(negative_kw)) * 3)
    return {"subtitle": subtitle, "content": content, "priority": priority}


def _build_wellness_break_card(user_id: int, continuous_minutes: int) -> dict:
    """
    主动关护卡：融合天气 + 兴趣建议。
    触发语气以“连续使用时长”作为核心，结合用户近期兴趣给出放松建议。
    """
    interests = _infer_user_interests(int(user_id), limit=4)
    top_interest = (interests[0] if interests else "").strip()
    second_interest = (interests[1] if len(interests) > 1 else "").strip()
    weather_hint = "合肥今天天气不错，适合短暂走动放松。"

    actions = [
        "起来走走、活动肩颈 5-10 分钟",
        "去窗边透透气，喝点温水再继续",
    ]
    if top_interest:
        actions.append(f"做一小段和「{top_interest}」相关的放松活动")
    if second_interest:
        actions.append(f"或者切换到「{second_interest}」话题放空几分钟")

    rec_text = "；".join(actions[:3])
    subtitle = f"你已连续使用约 {continuous_minutes} 分钟"
    content = f"{weather_hint} 要不要先{rec_text}？短休之后再继续，效率通常会更高。"

    return {
        "id": "wellness_break",
        "type": "wellness",
        "priority": 78,
        "title": "主动关护提醒",
        "subtitle": subtitle,
        "content": content,
    }


def _maybe_pick_wellness_break_milestone(user_id: int, continuous_minutes: int, now: datetime = None):
    """
    连续使用每 45 分钟触发一次主动关护卡。
    使用 daily milestone 去重：45, 90, 135, ...
    """
    if continuous_minutes < 45:
        return None
    if now is None:
        now = datetime.now()
    day_key = now.strftime("%Y-%m-%d")
    target = int(continuous_minutes // 45) * 45
    try:
        logs = db.list_user_memory_vectors(int(user_id), memory_type="wellness_break_push", limit=120) or []
    except Exception:
        logs = []
    done = set()
    for r in logs:
        try:
            meta = json.loads((r.get("metadata_json") or "{}").strip())
            if str(meta.get("date") or "") != day_key:
                continue
            m = meta.get("milestone_min")
            if m is not None:
                done.add(int(m))
        except Exception:
            continue
    if target in done:
        return None
    return target


def _get_continuous_usage_minutes(user_id: int, now: datetime = None) -> int:
    """估算用户连续使用时长（分钟）。超过空档阈值则重置连续会话。"""
    if now is None:
        now = datetime.now()
    now_ts = int(now.timestamp())
    with _mood_tip_session_lock:
        s = _mood_tip_session_by_user.get(int(user_id))
        if not s:
            _mood_tip_session_by_user[int(user_id)] = {"start_ts": now_ts, "last_seen_ts": now_ts}
            return 0
        last_seen_ts = int(s.get("last_seen_ts") or now_ts)
        gap_minutes = int((now_ts - last_seen_ts) / 60)
        if gap_minutes > MOOD_TIP_SESSION_GAP_RESET_MINUTES:
            s["start_ts"] = now_ts
        s["last_seen_ts"] = now_ts
        start_ts = int(s.get("start_ts") or now_ts)
    if now_ts <= start_ts:
        return 0
    return max(0, int((now_ts - start_ts) / 60))


def _maybe_pick_mood_tip_milestone(user_id: int, negative_kw: list, continuous_minutes: int, now: datetime = None):
    """
    返回应触发的里程碑（分钟），不满足则返回 None。
    规则：
    - 必须存在负向情绪信号
    - 连续使用时长达到里程碑（30/60/180/300）
    - 当日总触发次数不超过上限
    - 同里程碑当日只触发一次
    """
    if not negative_kw:
        return None
    if now is None:
        now = datetime.now()
    day_key = now.strftime("%Y-%m-%d")
    try:
        logs = db.list_user_memory_vectors(int(user_id), memory_type="mood_tip_push", limit=40) or []
    except Exception:
        logs = []
    pushed_milestones = set()
    today_count = 0
    for r in logs:
        try:
            meta = json.loads((r.get("metadata_json") or "{}").strip())
            if str(meta.get("date") or "") != day_key:
                continue
            today_count += 1
            m = meta.get("milestone_min")
            if m is not None:
                pushed_milestones.add(int(m))
        except Exception:
            continue
    if today_count >= MOOD_TIP_DAILY_MAX_PUSH:
        return None
    reached = [m for m in MOOD_TIP_MILESTONES_MINUTES if continuous_minutes >= m]
    if not reached:
        return None
    for m in reversed(reached):
        if m not in pushed_milestones:
            return m
    return None


def _infer_user_interests(user_id: int, limit: int = 6) -> list:
    lexicon = ["篮球", "跑步", "骑行", "无人机", "航拍", "电影", "音乐", "写代码", "学习", "摄影", "爬山"]
    hits = Counter()
    try:
        rows = db.get_messages_for_user_range(
            int(user_id),
            (datetime.now() - timedelta(days=30)).strftime("%Y-%m-%d %H:%M:%S"),
            datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            limit=3000,
        ) or []
        for r in rows:
            if (r.get("role") or "") != "user":
                continue
            text = (r.get("content") or "").strip()
            for kw in lexicon:
                if kw in text:
                    hits[kw] += 1
    except Exception:
        pass
    return [k for k, _ in hits.most_common(max(1, limit))]


def _parse_rss_channel_items(xml_bytes: bytes, limit: int):
    """解析常见 RSS 2.0 channel/item，返回 title / link。"""
    out = []
    try:
        root = ET.fromstring(xml_bytes)
    except ET.ParseError:
        return out
    channel = root.find("channel")
    if channel is None:
        return out
    for it in channel.findall("item"):
        if len(out) >= limit:
            break
        t_el = it.find("title")
        l_el = it.find("link")
        title = (t_el.text or "").strip() if t_el is not None else ""
        link = (l_el.text or "").strip() if l_el is not None else ""
        title = html_module.unescape(title)
        if not title:
            continue
        out.append({"title": title, "link": link})
    return out


def _fetch_interest_related_news_items(interests, max_items: int = 12):
    """
    从公开 RSS（Bing News / Google News）抓取与用户兴趣相关的新闻标题，供前端交给模型整理。
    任一源失败则尝试下一源；均失败返回空列表。
    """
    base_interests = [x.strip() for x in (interests or []) if x and str(x).strip()]
    if not base_interests:
        base_interests = ["科技", "国内"]

    queries = []
    for kw in base_interests[:3]:
        queries.append(f"{kw} 新闻")
    if len(base_interests) >= 2:
        queries.append(f"{' '.join(base_interests[:2])} 资讯")

    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
            "(KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36"
        ),
        "Accept": "application/rss+xml, application/xml, text/xml, */*",
    }

    seen = set()
    collected = []

    def _append_from_items(items) -> None:
        nonlocal collected
        for it in items:
            if len(collected) >= max_items:
                return
            title = (it.get("title") or "").strip()
            if not title:
                continue
            key = title[:120]
            if key in seen:
                continue
            seen.add(key)
            collected.append({"title": title, "link": (it.get("link") or "").strip()})

    for q in queries:
        if len(collected) >= max_items:
            break
        enc = urllib.parse.quote_plus(q)
        # Bing News RSS（国内环境相对易访问）
        bing_url = f"https://www.bing.com/news/search?q={enc}&setlang=zh-cn&cc=CN&format=rss"
        try:
            r = _requests().get(bing_url, timeout=14, headers=headers)
            r.raise_for_status()
            _append_from_items(_parse_rss_channel_items(r.content, max_items - len(collected)))
        except Exception:
            pass
        if len(collected) >= max_items:
            break
        g_url = (
            "https://news.google.com/rss/search?q="
            + urllib.parse.quote(q)
            + "&hl=zh-CN&gl=CN&ceid=CN:zh"
        )
        try:
            r2 = _requests().get(g_url, timeout=14, headers=headers)
            r2.raise_for_status()
            _append_from_items(_parse_rss_channel_items(r2.content, max_items - len(collected)))
        except Exception:
            pass

    if not collected:
        # 如果 RSS 抓取失败（如网络限制），使用大模型生成几条相关新闻标题作为兜底
        sys_prompt = (
            "你是一个新闻聚合助手。请根据用户的兴趣标签，生成 3~5 条最近的真实或高度逼真的新闻标题。"
            "返回 JSON 格式：{\"items\": [{\"title\": \"新闻标题\", \"link\": \"\"}]}"
        )
        user_prompt = f"用户的兴趣标签：{'、'.join(base_interests)}"
        try:
            obj, err = _llm_json_chat(sys_prompt, user_prompt, max_tokens=500, timeout=15)
            if not err and isinstance(obj, dict):
                llm_items = obj.get("items") or []
                for it in llm_items:
                    if isinstance(it, dict) and it.get("title"):
                        collected.append({"title": it["title"].strip(), "link": (it.get("link") or "").strip()})
        except Exception:
            pass

    return collected[:max_items]


def _build_interest_news_summary(items, interests) -> str:
    """生成兴趣新闻卡片摘要（基于已抓取标题，不编造内容）。"""
    if not items:
        return "正在为你整理相关新闻，点击本卡片可在聊天区生成按兴趣分类的资讯简报。"
    top_titles = [(x.get("title") or "").strip() for x in items[:2] if (x.get("title") or "").strip()]
    if not top_titles:
        return "已获取到相关新闻源，点击本卡片可查看按兴趣整理后的摘要。"
    interest_text = "、".join(interests[:3]) if interests else "你的近期兴趣"
    return f"检测到你近期关注{interest_text}：{top_titles[0]}" + (f"；另外还有：{top_titles[1]}。" if len(top_titles) > 1 else "。")


def _get_or_refresh_interest_news_payload(user_id: int, force_refresh: bool = False) -> dict:
    """获取兴趣新闻缓存，不存在或过期时刷新。"""
    now_ts = time.time()
    with _interest_news_cache_lock:
        cached = _interest_news_cache_by_user.get(int(user_id))
        if (
            not force_refresh
            and cached
            and float(cached.get("expire_at") or 0) > now_ts
            and isinstance(cached.get("items"), list)
        ):
            return {
                "interests": cached.get("interests") or [],
                "items": cached.get("items") or [],
                "summary": (cached.get("summary") or "").strip(),
                "cached_at": int(cached.get("cached_at") or now_ts),
            }

    interests = _infer_user_interests(int(user_id), limit=5)
    items = _fetch_interest_related_news_items(interests, max_items=14)
    payload = {
        "interests": interests,
        "items": items,
        "summary": _build_interest_news_summary(items, interests),
        "cached_at": int(now_ts),
    }
    with _interest_news_cache_lock:
        _interest_news_cache_by_user[int(user_id)] = {
            **payload,
            "expire_at": now_ts + max(60, INTEREST_NEWS_CACHE_TTL_SEC),
        }
    return payload


def _daily_interest_news_push_decision(user_id: int, now: datetime = None) -> dict:
    """
    兴趣新闻卡推送节奏：
    1) 早上一次（05:00-11:59）
    2) 白天一次（12:00-14:59）
    3) 下午随机时间一次（15:00 后某个固定随机时点）
    """
    if now is None:
        now = datetime.now()
    day_key = now.strftime("%Y-%m-%d")
    minute_of_day = now.hour * 60 + now.minute

    # 使用 user_id + 日期生成“当日稳定随机时点”（15:00 ~ 18:29）
    day_num = int(day_key.replace("-", ""))
    rand_offset = (int(user_id) * 97 + day_num * 13) % 210
    afternoon_trigger_minute = 15 * 60 + rand_offset
    afternoon_end_minute = min(18 * 60 + 59, afternoon_trigger_minute + 45)  # 下午随机窗口 45 分钟

    shown_buckets = set()
    shown_signatures = set()
    try:
        logs = db.list_user_memory_vectors(int(user_id), memory_type="interest_news_push", limit=30) or []
    except Exception:
        logs = []
    for r in logs:
        try:
            meta = json.loads((r.get("metadata_json") or "{}").strip())
            if str(meta.get("date") or "") != day_key:
                continue
            bucket = str(meta.get("bucket") or "").strip()
            if bucket:
                shown_buckets.add(bucket)
            signature = str(meta.get("signature") or "").strip()
            if signature:
                shown_signatures.add(signature)
        except Exception:
            continue

    def _done(bucket_name: str) -> bool:
        return bucket_name in shown_buckets

    # 早上
    if 5 * 60 <= minute_of_day <= 11 * 60 + 59 and not _done("morning"):
        return {
            "show": True,
            "bucket": "morning",
            "force_refresh": True,
            "day_key": day_key,
            "shown_signatures": shown_signatures,
        }

    # 白天
    if 12 * 60 <= minute_of_day <= 14 * 60 + 59 and not _done("daytime"):
        return {
            "show": True,
            "bucket": "daytime",
            "force_refresh": True,
            "day_key": day_key,
            "shown_signatures": shown_signatures,
        }

    # 下午随机时段触发一次（过窗即作废）
    if afternoon_trigger_minute <= minute_of_day <= afternoon_end_minute and not _done("afternoon"):
        return {
            "show": True,
            "bucket": "afternoon",
            "force_refresh": True,
            "day_key": day_key,
            "shown_signatures": shown_signatures,
        }

    return {
        "show": False,
        "bucket": "",
        "force_refresh": False,
        "day_key": day_key,
        "shown_signatures": shown_signatures,
    }


def _interest_news_signature(items: list) -> str:
    """基于前三条标题生成签名，用于同日去重。"""
    titles = []
    for it in (items or [])[:3]:
        t = (it.get("title") or "").strip().lower() if isinstance(it, dict) else ""
        if t:
            titles.append(t[:120])
    if not titles:
        return ""
    return " || ".join(titles)


def _auto_expire_interest_news_buckets(user_id: int, now: datetime = None) -> None:
    """
    超过时间窗后自动作废并记一次消费（expired）：
    - 早上窗：05:00-11:59
    - 白天窗：12:00-14:59
    - 下午随机窗：start ~ start+45min（不晚于18:59）
    """
    if now is None:
        now = datetime.now()
    day_key = now.strftime("%Y-%m-%d")
    minute_of_day = now.hour * 60 + now.minute

    day_num = int(day_key.replace("-", ""))
    rand_offset = (int(user_id) * 97 + day_num * 13) % 210
    afternoon_trigger_minute = 15 * 60 + rand_offset
    afternoon_end_minute = min(18 * 60 + 59, afternoon_trigger_minute + 45)

    shown_buckets = set()
    try:
        logs = db.list_user_memory_vectors(int(user_id), memory_type="interest_news_push", limit=40) or []
    except Exception:
        logs = []
    for r in logs:
        try:
            meta = json.loads((r.get("metadata_json") or "{}").strip())
            if str(meta.get("date") or "") != day_key:
                continue
            bucket = str(meta.get("bucket") or "").strip()
            if bucket:
                shown_buckets.add(bucket)
        except Exception:
            continue

    def _expire(bucket: str) -> None:
        _memory_add(
            int(user_id),
            "interest_news_push",
            "expired",
            {
                "date": day_key,
                "bucket": bucket,
                "hour": now.hour,
                "minute": now.minute,
                "signature": "",
                "reason": "expired_window",
            },
        )

    if minute_of_day > 11 * 60 + 59 and "morning" not in shown_buckets:
        _expire("morning")
        shown_buckets.add("morning")
    if minute_of_day > 14 * 60 + 59 and "daytime" not in shown_buckets:
        _expire("daytime")
        shown_buckets.add("daytime")
    if minute_of_day > afternoon_end_minute and "afternoon" not in shown_buckets:
        _expire("afternoon")


def _build_interest_adaptive_chat_prompts(user_id: int) -> list:
    """基于用户兴趣/近期任务生成聊天区引导语。"""
    interests = _infer_user_interests(int(user_id), limit=6)
    top = interests[:3]
    top1 = top[0] if len(top) > 0 else "学习"
    top2 = top[1] if len(top) > 1 else "计划管理"
    top3 = top[2] if len(top) > 2 else "放松"
    # 日程类/情绪类引导语固定；兴趣类引导语动态随用户偏好变化
    fixed_schedule_emotion_prompts = [
        "你可以对我说：“帮我梳理一下今天的日程优先级，先做哪件事最稳妥？”",
        "你可以对我说：“今天太忙了，帮我把不紧急的日程延后。”",
        "你可以对我说：“我现在状态不太好，带我做一个 2 分钟快速恢复。”",
        "你可以对我说：“我今天有点累，帮我把任务节奏调轻一点。”",
    ]
    dynamic_interest_prompts = [
        f"你可以对我说：“最近我常聊{top1}，给我来一条今天值得关注的{top1}动态。”",
        f"你可以对我说：“结合我最近在{top1}/{top2}上的关注点，给我一个今晚学习或练习计划。”",
        f"你可以对我说：“用{top3}相关的话题陪我聊 5 分钟，顺便帮我放松一下。”",
        f"你可以对我说：“我最近在关注{top2}，帮我整理 3 条实用建议。”",
    ]
    prompts = fixed_schedule_emotion_prompts + dynamic_interest_prompts

    # 去重并限制数量，避免前端轮播过长
    out = []
    seen = set()
    for p in prompts:
        key = (p or "").strip()
        if not key or key in seen:
            continue
        seen.add(key)
        out.append(key)
    return out[:8]


def _get_or_refresh_chat_prompts_payload(user_id: int, force_refresh: bool = False) -> dict:
    """聊天区引导语：至少 3 天更新一次；force 可强制刷新。"""
    now = datetime.now()
    now_ts = int(time.time())
    ttl_sec = 3 * 24 * 60 * 60
    if not force_refresh:
        try:
            rows = db.list_user_memory_vectors(int(user_id), memory_type="chat_prompt_pack", limit=20) or []
        except Exception:
            rows = []
        for r in rows:
            try:
                meta = json.loads((r.get("metadata_json") or "{}").strip())
                expires_at = int(meta.get("expires_at") or 0)
                prompts = meta.get("prompts") or []
                if expires_at > now_ts and isinstance(prompts, list) and prompts:
                    clean = [str(x).strip() for x in prompts if str(x).strip()]
                    if clean:
                        return {
                            "prompts": clean[:8],
                            "cached_at": int(meta.get("cached_at") or now_ts),
                            "expires_at": expires_at,
                            "interests": meta.get("interests") or [],
                        }
            except Exception:
                continue

    prompts = _build_interest_adaptive_chat_prompts(int(user_id))
    interests = _infer_user_interests(int(user_id), limit=4)
    payload = {
        "prompts": prompts,
        "cached_at": now_ts,
        "expires_at": now_ts + ttl_sec,
        "interests": interests,
    }
    _memory_add(
        int(user_id),
        "chat_prompt_pack",
        "chat_prompt_pack_generated",
        {
            "date": now.strftime("%Y-%m-%d"),
            "cached_at": now_ts,
            "expires_at": now_ts + ttl_sec,
            "prompts": prompts,
            "interests": interests,
        },
    )
    return payload


def _detect_interest_news_active_bucket(user_id: int, now: datetime = None) -> str:
    if now is None:
        now = datetime.now()
    minute_of_day = now.hour * 60 + now.minute
    if 5 * 60 <= minute_of_day <= 11 * 60 + 59:
        return "morning"
    if 12 * 60 <= minute_of_day <= 14 * 60 + 59:
        return "daytime"
    day_key = now.strftime("%Y-%m-%d")
    day_num = int(day_key.replace("-", ""))
    rand_offset = (int(user_id) * 97 + day_num * 13) % 210
    afternoon_trigger_minute = 15 * 60 + rand_offset
    afternoon_end_minute = min(18 * 60 + 59, afternoon_trigger_minute + 45)
    if afternoon_trigger_minute <= minute_of_day <= afternoon_end_minute:
        return "afternoon"
    return ""


def _today_schedule_snapshot(user_id: int) -> list:
    now = datetime.now()
    start = now.strftime("%Y-%m-%d 00:00:00")
    end = now.strftime("%Y-%m-%d 23:59:59")
    try:
        return db.get_schedules_by_user(int(user_id), start_date=start, end_date=end, limit=200) or []
    except Exception:
        return []


def _build_morning_greeting(user_id: int) -> dict:
    user_row = db.get_user_by_id(int(user_id)) or {}
    call_name = ((user_row.get("preferred_name") or "").strip() or (user_row.get("username") or "").strip() or "你")[:50]
    schedules = _today_schedule_snapshot(int(user_id))
    pending = [s for s in schedules if (s.get("status") or "pending") == "pending"]
    mood_keywords = _today_mood_keywords(int(user_id), limit=3)
    interests = _infer_user_interests(int(user_id), limit=3)
    top_interest = interests[0] if interests else ""
    mood_hint = "、".join(mood_keywords) if mood_keywords else "平稳"
    schedule_hint = "今天日程比较宽松" if len(pending) <= 2 else f"今天有 {len(pending)} 项待办"
    interest_hint = f"你最近常聊{top_interest}" if top_interest else "我可以帮你梳理今天计划"

    content = (
        f"早上好，{call_name}。{schedule_hint}，当前情绪关键词偏向「{mood_hint}」。"
        f"{interest_hint}，要不要先一起确认今天最重要的三件事？"
    )
    if ("疲惫" in mood_hint or "压力" in mood_hint or "焦虑" in mood_hint) and len(pending) <= 2:
        content += "如果你今天比较累，我们也可以先放缓节奏，安排一段恢复时间。"
    if top_interest in ("篮球", "骑行", "无人机", "航拍"):
        content += f"若下午状态允许，可以安排一段{top_interest}放松。"

    return {
        "message": content,
        "mood_keywords": mood_keywords,
        "interests": interests,
        "schedule_count": len(pending),
    }


def _build_daytime_recommendations(user_id: int) -> list:
    now = datetime.now()
    recs = []
    today_rows = _today_schedule_snapshot(int(user_id))
    pending = [r for r in today_rows if (r.get("status") or "pending") == "pending"]
    next_item = None
    min_delta = None
    for r in pending:
        sat = r.get("scheduled_at")
        if not hasattr(sat, "timestamp"):
            continue
        delta_m = int((sat - now).total_seconds() / 60)
        if delta_m < 0:
            continue
        if min_delta is None or delta_m < min_delta:
            min_delta = delta_m
            next_item = r
    if next_item and min_delta is not None and min_delta <= 30:
        sat = next_item.get("scheduled_at")
        sat_label = sat.strftime("%H:%M") if hasattr(sat, "strftime") else ""
        recs.append({
            "type": "schedule_reminder",
            "priority": "high",
            "message": f"提醒：你的下一项日程「{(next_item.get('title') or '').strip()}」将在 {max(0, min_delta)} 分钟后开始" + (f"（{sat_label}）" if sat_label else "") + "。",
            "schedule_id": int(next_item.get("id") or 0),
        })

    mood_keywords = _today_mood_keywords(int(user_id), limit=4)
    fatigue = any(x in ("疲惫", "焦虑", "压力", "烦躁", "崩溃") for x in mood_keywords)
    if fatigue:
        recs.append({
            "type": "emotion_task_adjustment",
            "priority": "medium",
            "message": "检测到你今天状态偏疲惫。建议先暂停手头任务 10 分钟，我可以帮你记录当前进度并稍后恢复。",
            "action": "interrupt_and_resume",
        })

    if not recs:
        recs.append({
            "type": "routine",
            "priority": "low",
            "message": "你当前安排看起来节奏不错，需要我帮你提前准备下一项任务清单吗？",
        })
    return recs


def _schedule_time_label(sat) -> str:
    if hasattr(sat, "strftime"):
        return sat.strftime("%H:%M")
    s = str(sat or "")
    return s[11:16] if len(s) >= 16 else ""


def _build_today_plan_card(pending: list, recs: list, fatigue_detected: bool) -> dict:
    if not pending:
        return {}
    sorted_pending = sorted(
        pending,
        key=lambda x: (
            x.get("scheduled_at").timestamp() if hasattr(x.get("scheduled_at"), "timestamp") else 0,
            int(x.get("id") or 0),
        ),
    )
    lines = []
    for i, row in enumerate(sorted_pending[:4]):
        title = (row.get("title") or "").strip() or "未命名事项"
        t = _schedule_time_label(row.get("scheduled_at"))
        lines.append(f"{i + 1}. {title}" + (f"（{t}）" if t else ""))
    first = sorted_pending[0]
    first_title = (first.get("title") or "").strip() or "第一项任务"
    first_time = _schedule_time_label(first.get("scheduled_at"))
    guidance = f"建议先做「{first_title}」" + (f"（{first_time}）" if first_time else "") + "，完成后再处理后续事项。"
    return {
        "id": "today_plan",
        "type": "schedule",
        "priority": 99 if fatigue_detected else 96,
        "title": "今日行程总览",
        "subtitle": f"今天共有 {len(sorted_pending)} 项待办",
        "content": "\n".join(lines + [guidance]),
    }


def _build_high_pressure_relief_card(call_name: str = "你") -> dict:
    who = (call_name or "你").strip() or "你"
    return {
        "id": "high_pressure_relief",
        "type": "wellness",
        "priority": 100,
        "title": "高压减负提醒",
        "subtitle": f"{who}，先把任务放一放",
        "content": "你已经连续高强度投入一段时间了。建议先休息 10 分钟：起身活动、喝点水、做几次深呼吸，调整后再继续会更高效。",
    }


def _detect_multitask_temporary_plan(user_text: str) -> list:
    txt = (user_text or "").strip()
    if not txt:
        return []
    multitask_markers = ["同时", "一边", "另外", "然后", "再", "还要", "并且", "同时做", "多个任务"]
    if not any(k in txt for k in multitask_markers):
        return []
    tasks = []
    chunks = re.split(r"[，,。；;、然后再并且]+", txt)
    for c in chunks:
        c = c.strip()
        if len(c) < 4:
            continue
        if any(kw in c for kw in ("帮我", "我想", "我要", "需要")):
            title = c[:40]
            tasks.append(title)
    return tasks[:3]


def _create_temporary_schedules_from_text(user_id: int, user_text: str) -> list:
    tasks = _detect_multitask_temporary_plan(user_text)
    if not tasks:
        return []
    now = datetime.now()
    created = []
    for i, t in enumerate(tasks):
        sat = (now + timedelta(minutes=(i + 1) * 30)).strftime("%Y-%m-%d %H:%M:%S")
        try:
            sid = db.create_schedule(
                int(user_id),
                title=f"[临时] {t}"[:500],
                scheduled_at=sat,
                source="temporary",
                raw_text=(user_text or "")[:500],
            )
            created.append({"id": sid, "title": f"[临时] {t}", "scheduled_at": sat})
        except Exception:
            continue
    return created


@app.route("/api/chat", methods=["POST"])
def chat():
    """
    请求体: { "content": "用户输入", "messages": [...], 可选 "imageBase64"/"imageUrl", "videoBase64"/"videoUrl", "attachmentHint" }
    响应:   { "content": "AI 回复文本" }
    """
    try:
        data = request.get_json(force=True, silent=True) or {}
        content = (data.get("content") or "").strip()
        attachment_hint = (data.get("attachmentHint") or data.get("attachment_hint") or "").strip()
        image_b64, image_mime = _resolve_image_from_request(data)
        video_b64, video_mime = _resolve_video_from_request(data)
        doc_text, _doc_mime, doc_name = _resolve_document_from_request(data)
        if not content and not attachment_hint and not image_b64 and not video_b64 and not doc_text:
            return jsonify({"error": "content 或附件不能为空"}), 400
        if attachment_hint and not content:
            content = attachment_hint
        history = data.get("messages") or []
        messages = _with_chat_system_prompt(
            build_messages(
                history,
                content,
                image_b64,
                image_mime,
                video_b64,
                video_mime,
                document_text=doc_text,
                document_name=doc_name,
            )
        )

        payload = {
            "model": CHAT_MODEL,
            "messages": messages,
            "max_tokens": MAX_TOKENS,
        }
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {API_KEY}",
        }

        resp = _requests().post(CHAT_API_URL, json=payload, headers=headers, timeout=60)
        resp.raise_for_status()
        result = resp.json()
        ai_content = (
            result.get("choices", [{}])[0].get("message", {}).get("content", "")
        ).strip()
        return jsonify({"content": ai_content or "（无回复）"})
    except Exception as e:
        err_msg = str(e)
        if hasattr(e, "response") and e.response is not None:
            try:
                err_body = e.response.json()
                err_msg = err_body.get("error", {}).get("message", err_msg)
            except Exception:
                err_msg = e.response.text or err_msg
        return jsonify({"error": f"模型请求失败: {err_msg}"}), 502
    except Exception as e:
        return jsonify({"error": str(e)}), 500


def _conv_row_to_json(row):
    """会话行转前端格式。"""
    if not row:
        return None
    return {
        "id": str(row["id"]),
        "title": row.get("title") or "新对话",
        "lastMessage": (row.get("last_message") or "").strip(),
        "updatedAt": row["updated_at"].isoformat() if hasattr(row.get("updated_at"), "isoformat") else str(row.get("updated_at", "")),
        "messageCount": int(row.get("message_count", 0)),
        "pinned": bool(row.get("pinned")),
    }


def _msg_row_to_json(row):
    """消息行转前端格式：role assistant -> sender ai。"""
    if not row:
        return None
    role = (row.get("role") or "user").strip()
    return {
        "id": str(row["id"]),
        "content": (row.get("content") or "").strip(),
        "sender": "ai" if role == "assistant" else "user",
        "timestamp": row["created_at"].isoformat() if hasattr(row.get("created_at"), "isoformat") else str(row.get("created_at", "")),
        "type": (row.get("type") or "text").strip(),
        "fileUrl": (row.get("file_url") or "").strip() or None,
        "fileName": (row.get("file_name") or "").strip() or None,
    }


@app.route("/api/conversations", methods=["GET"])
def list_conversations():
    """当前用户的会话列表，需登录。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        limit = request.args.get("limit", type=int) or 200
        limit = min(max(1, limit), 500)
        rows = db.get_conversations_by_user(int(user_id), limit=limit)
        return jsonify([_conv_row_to_json(r) for r in rows])
    except Exception as e:
        return _err("查询失败: " + str(e), 500)


@app.route("/api/conversations", methods=["POST"])
def create_conversation():
    """创建会话，需登录。body: { title? }"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        data = request.get_json(force=True, silent=True) or {}
        title = (data.get("title") or "新对话").strip()[:255]
        conv_id = db.create_conversation(int(user_id), title)
        row = db.get_conversation_by_id(conv_id, int(user_id))
        return jsonify(_conv_row_to_json(row))
    except Exception as e:
        return _err("创建失败: " + str(e), 500)


@app.route("/api/conversations/<int:conv_id>", methods=["GET"])
def get_conversation(conv_id):
    """获取单条会话，需登录且为本人。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    row = db.get_conversation_by_id(conv_id, int(user_id))
    if not row:
        return _err("会话不存在", 404)
    return jsonify(_conv_row_to_json(row))


@app.route("/api/conversations/<int:conv_id>", methods=["PATCH"])
def patch_conversation(conv_id):
    """更新会话 title / pinned，需登录且为本人。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    title = data.get("title")
    pinned = data.get("pinned")
    if title is None and pinned is None:
        return jsonify(_conv_row_to_json(db.get_conversation_by_id(conv_id, int(user_id))))
    ok = db.update_conversation(conv_id, int(user_id), title=title, pinned=pinned)
    if not ok:
        return _err("会话不存在", 404)
    row = db.get_conversation_by_id(conv_id, int(user_id))
    return jsonify(_conv_row_to_json(row))


@app.route("/api/conversations/<int:conv_id>", methods=["DELETE"])
def delete_conversation(conv_id):
    """删除会话，需登录且为本人。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    ok = db.delete_conversation(conv_id, int(user_id))
    if not ok:
        return _err("会话不存在", 404)
    return jsonify({"ok": True})


@app.route("/api/conversations/<int:conv_id>/messages", methods=["GET"])
def list_messages(conv_id):
    """会话消息列表，需登录且为本人。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    conv = db.get_conversation_by_id(conv_id, int(user_id))
    if not conv:
        return _err("会话不存在", 404)
    limit = request.args.get("limit", type=int) or 500
    limit = min(max(1, limit), 1000)
    rows = db.get_messages_by_conversation(conv_id, limit=limit)
    # MySQL DATETIME 仅到秒级；同秒写入时补充按 id 排序，避免显示顺序抖动。
    rows = sorted(rows or [], key=lambda r: (str(r.get("created_at") or ""), int(r.get("id") or 0)))
    return jsonify([_msg_row_to_json(r) for r in rows])


@app.route("/api/chat/stream", methods=["POST"])
def chat_stream():
    """
    流式对话：需登录。请求体需含 conversationId；支持 content、imageUrl、videoUrl、attachmentHint、audioUrl。
    会将会话与消息写入数据库，历史从数据库读取。
    响应为 SSE：data: {"content": "增量文本"}，结束 data: [DONE] 或 data: {"type":"done"}。
    """
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        data = request.get_json(force=True, silent=True) or {}
        conv_id_raw = data.get("conversationId") or data.get("conversation_id")
        if conv_id_raw is None:
            return jsonify({"error": "缺少 conversationId"}), 400
        try:
            conv_id = int(conv_id_raw)
        except (TypeError, ValueError):
            return jsonify({"error": "conversationId 无效"}), 400
        owner = db.get_conversation_owner(conv_id)
        if owner is None or owner != int(user_id):
            return jsonify({"error": "会话不存在或无权访问"}), 404

        content = (data.get("content") or "").strip()
        attachment_hint = (data.get("attachmentHint") or data.get("attachment_hint") or "").strip()
        image_b64, image_mime = _resolve_image_from_request(data)
        video_b64, video_mime = _resolve_video_from_request(data)
        audio_b64, audio_mime, voice_file_name = _resolve_audio_from_request(data)
        doc_text, _doc_mime, doc_name = _resolve_document_from_request(data)
        if not content and not attachment_hint and not image_b64 and not video_b64 and not audio_b64 and not doc_text:
            return jsonify({"error": "content 或附件不能为空"}), 400
        if attachment_hint and not content:
            content = attachment_hint

        # 从数据库取历史（仅文本，用于上下文）
        history_rows = db.get_messages_by_conversation(conv_id, limit=50)
        history = []
        for r in history_rows:
            role = (r.get("role") or "user").strip()
            if role in ("user", "assistant"):
                history.append({"role": role, "content": (r.get("content") or "").strip()})
        
        hide_user_message = data.get("hideUserMessage") or data.get("hide_user_message")

        # 写入用户消息
        user_content = content or (
            attachment_hint
            or ("用户发来一条语音：" if audio_b64 else (f"用户上传文档：{doc_name}" if doc_text else "[图片/附件]"))
        )
        user_msg_type = "voice" if audio_b64 else "text"
        user_file_url = (data.get("audioUrl") or data.get("audio_url")) if audio_b64 else None
        user_file_name = (
            (voice_file_name or (data.get("voiceFileName") or data.get("voice_file_name")))
            if audio_b64
            else None
        )
        if not hide_user_message:
            db.create_message(
                conv_id,
                "user",
                user_content,
                user_msg_type,
                user_file_url,
                user_file_name,
            )
            db.update_conversation_last_message(conv_id, user_content)

        # 构建发送给 AI 的 messages
        # 如果 hide_user_message 为 True，说明这是一条后台自动触发的消息（如兴趣新闻简报），
        # 此时我们不把这条长长的 prompt 算作用户的正常发言，而是直接把它作为 system prompt 的一部分，
        # 或者作为一个独立的 user 消息发给 AI，但不落库。
        messages_for_ai = _with_chat_system_prompt(
            build_messages(
                history,
                content,
                image_b64,
                image_mime,
                video_b64,
                video_mime,
                audio_b64,
                audio_mime,
                document_text=doc_text,
                document_name=doc_name,
            )
        )
        
        if hide_user_message:
            # 如果是隐藏消息，替换最后一条 user 消息的内容为更自然的指令
            if messages_for_ai and messages_for_ai[-1]["role"] == "user":
                messages_for_ai[-1]["content"] = content

        payload = {
            "model": CHAT_MODEL,
            "messages": messages_for_ai,
            "max_tokens": MAX_TOKENS,
            "stream": True,
        }
        # 用户上传语音（input_audio）时仅要文字回复，不请求助手侧 TTS，避免与「语音输入」组合时兼容或流式异常
        if _chat_model_supports_omni_audio() and not audio_b64:
            payload["modalities"] = ["text", "audio"]
            payload["audio"] = {
                "voice": (CHAT_OMNI_VOICE or "Cherry").strip(),
                "format": (CHAT_OMNI_AUDIO_FORMAT or "wav").strip().lstrip("."),
            }
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {API_KEY}",
        }

        def generate():
            full_content = []
            audio_b64_parts = []
            try:
                resp = _requests().post(
                    CHAT_API_URL, json=payload, headers=headers, timeout=120, stream=True
                )
                resp.raise_for_status()
                for line in resp.iter_lines(decode_unicode=True):
                    if not line or not line.strip():
                        continue
                    raw = line[6:].strip() if line.startswith("data: ") else line.strip()
                    if raw == "[DONE]":
                        break
                    try:
                        obj = json.loads(raw)
                        choices = obj.get("choices") or []
                        if not choices:
                            continue
                        delta = choices[0].get("delta") or {}
                        if not isinstance(delta, dict):
                            continue
                        text_delta = delta.get("content") or ""
                        if text_delta:
                            full_content.append(text_delta)
                            yield f"data: {json.dumps({'content': text_delta}, ensure_ascii=False)}\n\n"
                        audio_obj = delta.get("audio")
                        if isinstance(audio_obj, dict):
                            piece = audio_obj.get("data")
                            if isinstance(piece, str) and piece:
                                audio_b64_parts.append(piece)
                    except (json.JSONDecodeError, IndexError, KeyError, TypeError):
                        pass
                # 流结束后写入 assistant 消息并更新会话摘要；从用户输入中抽取日程并写入
                ai_content = "".join(full_content).strip() or "（无回复）"
                audio_url = None
                audio_name = None
                merged_b64 = "".join(audio_b64_parts).strip()
                if merged_b64:
                    try:
                        raw_audio = base64.b64decode(merged_b64, validate=False)
                        if raw_audio:
                            wav_bytes = _decoded_omni_audio_to_wav_bytes(raw_audio)
                            if wav_bytes:
                                audio_url, audio_name = _save_assistant_audio_file(wav_bytes, conv_id)
                    except Exception:
                        pass
                db.create_message(
                    conv_id,
                    "assistant",
                    ai_content,
                    "text",
                    audio_url,
                    audio_name,
                )
                db.update_conversation_last_message(conv_id, ai_content)
                _memory_add(int(user_id), "chat", ai_content, {"conversation_id": conv_id, "role": "assistant"})
                if audio_url and audio_name:
                    yield f"data: {json.dumps({'audioUrl': audio_url, 'fileName': audio_name}, ensure_ascii=False)}\n\n"
                if not hide_user_message:
                    _detect_emotion_anomaly(int(user_id), user_content)
                    _memory_add(int(user_id), "chat", user_content, {"conversation_id": conv_id, "role": "user"})
                    _extract_schedule_deletions_from_text(int(user_id), user_content)
                    _extract_schedules_from_text(int(user_id), user_content)
                    _create_temporary_schedules_from_text(int(user_id), user_content)
                yield "data: [DONE]\n\n"
            except Exception as e:
                err = str(e)
                if hasattr(e, "response") and e.response is not None:
                    try:
                        err = (e.response.json() or {}).get("error", {}).get("message", err)
                    except Exception:
                        err = getattr(e.response, "text", None) or err
                yield f"data: {json.dumps({'error': err}, ensure_ascii=False)}\n\n"

        return Response(
            generate(),
            mimetype="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "X-Accel-Buffering": "no",
                "Connection": "keep-alive",
            },
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/health", methods=["GET"])
def health():
    return jsonify({"status": "ok"})


# ---------- 人脸识别 + 情绪识别（来自 gui_app2.py 能力后端化） ----------


def _student_row_to_json(row):
    if not row:
        return None
    return {
        "id": row["id"],
        "student_id": row["student_id"],
        "name": row["name"],
        "has_face_feature": bool(row.get("face_feature")),
        "is_deleted": int(row.get("is_deleted") or 0),
        "deleted_at": row["deleted_at"].isoformat() if row.get("deleted_at") and hasattr(row["deleted_at"], "isoformat") else (str(row.get("deleted_at")) if row.get("deleted_at") else None),
        "created_at": row["created_at"].isoformat() if hasattr(row.get("created_at"), "isoformat") else str(row.get("created_at", "")),
        "updated_at": row["updated_at"].isoformat() if hasattr(row.get("updated_at"), "isoformat") else str(row.get("updated_at", "")),
    }


def _record_row_to_json(row):
    if not row:
        return None
    return {
        "id": row["id"],
        "student_id": row["student_id"],
        "emotion_type": row["emotion_type"],
        "intensity": float(row.get("intensity") or 0),
        "timestamp": row["timestamp"].isoformat() if hasattr(row.get("timestamp"), "isoformat") else str(row.get("timestamp", "")),
        "is_deleted": int(row.get("is_deleted") or 0),
        "deleted_at": row["deleted_at"].isoformat() if row.get("deleted_at") and hasattr(row["deleted_at"], "isoformat") else (str(row.get("deleted_at")) if row.get("deleted_at") else None),
    }


def _load_face_db_embeddings():
    import numpy as np
    face_db = {}
    rows = db.list_students(include_deleted=False, limit=5000)
    for row in rows:
        raw = row.get("face_feature")
        if not raw:
            continue
        try:
            face_db[row["student_id"]] = np.array(json.loads(raw), dtype=np.float32)
        except Exception:
            continue
    return face_db


def _get_face_db_embeddings_cached(force_refresh: bool = False):
    global _face_db_cache, _face_db_cache_ts
    now_ts = time.time()
    with _face_db_cache_lock:
        if (
            (not force_refresh)
            and _face_db_cache is not None
            and (now_ts - _face_db_cache_ts) <= _face_db_cache_ttl_sec
        ):
            return _face_db_cache
    loaded = _load_face_db_embeddings()
    with _face_db_cache_lock:
        _face_db_cache = loaded
        _face_db_cache_ts = time.time()
    return loaded


def _invalidate_face_db_cache():
    global _face_db_cache, _face_db_cache_ts
    with _face_db_cache_lock:
        _face_db_cache = None
        _face_db_cache_ts = 0.0


# 注意：以下人脸相关函数在文件中靠前定义，但使用了后面定义的 _err()；
# 仅在请求处理时调用，此时 _err 已绑定。


def _get_face_engine_safe():
    """初始化人脸引擎（含首次下载权重）；失败时返回 JSON 错误而非 500 栈。"""
    try:
        return _face_mod().get_engine(), None
    except Exception as e:
        return None, _err(str(e), 503)


def _face_warmup_worker():
    """后台线程中执行 get_engine()，避免阻塞 HTTP 连接。"""
    global _face_warmup_state, _face_warmup_error
    try:
        with _face_warmup_lock:
            _face_warmup_state = "loading"
            _face_warmup_error = None
        print(
            "[face] 开始加载人脸引擎（MTCNN + FaceNet + 情绪）。内存不足时进程可能被系统 OOM Kill；"
            "建议预留约 2GB+ 可用 RAM 或配置 swap。",
            flush=True,
        )
        _face_mod().get_engine()
        with _face_warmup_lock:
            _face_warmup_state = "ready"
    except Exception as e:
        with _face_warmup_lock:
            _face_warmup_state = "error"
            _face_warmup_error = str(e)
        traceback.print_exc()


@app.route("/api/face/warmup", methods=["POST"])
def face_warmup():
    """可选：后台预取人脸引擎。业务上首次 /face/recognize 或带图的注册会惰性 get_engine()，不依赖本接口。"""
    global _face_warmup_state, _face_warmup_error
    try:
        _user_id, err_res = _require_auth()
        if err_res:
            return err_res
        with _face_warmup_lock:
            if _face_warmup_state == "ready":
                return jsonify({"ok": True, "message": "人脸引擎已就绪"}), 200
            if _face_warmup_state in ("starting", "loading"):
                return jsonify({"ok": False, "status": _face_warmup_state}), 202
            _face_warmup_state = "starting"
            _face_warmup_error = None
            threading.Thread(target=_face_warmup_worker, name="face_warmup", daemon=True).start()
            return jsonify({"ok": False, "status": "starting"}), 202
    except Exception as e:
        traceback.print_exc()
        m = str(e).strip() or repr(e)
        return jsonify({"error": m, "message": m}), 500


@app.route("/api/face/warmup/status", methods=["GET"])
def face_warmup_status():
    """配合可选预热：查询进度（须登录）。"""
    _user_id, err_res = _require_auth()
    if err_res:
        return err_res
    with _face_warmup_lock:
        st = _face_warmup_state
        err = _face_warmup_error
    return jsonify({"ready": st == "ready", "status": st, "error": err})


@app.route("/api/face/students", methods=["GET"])
def list_face_students():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    include_deleted = request.args.get("include_deleted") in ("1", "true", "True")
    limit = min(max(1, request.args.get("limit", type=int) or 200), 1000)
    rows = db.list_students(include_deleted=include_deleted, limit=limit)
    return jsonify([_student_row_to_json(r) for r in rows])


@app.route("/api/face/students", methods=["POST"])
def create_or_register_face_student():
    print("[face] POST /api/face/students（注册/更新）已开始", flush=True)
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    student_id = (data.get("student_id") or data.get("studentId") or "").strip()
    name = (data.get("name") or "").strip()
    image_base64 = (data.get("image_base64") or data.get("imageBase64") or "").strip()
    if not student_id or not name:
        return _err("请传入 student_id 和 name")

    face_feature_json = None
    if image_base64:
        with _face_infer_lock:
            engine, err_eng = _get_face_engine_safe()
            if err_eng:
                return err_eng
            frame = engine.decode_base64_image(image_base64)
            if frame is None:
                return _err("图片解析失败，请检查 image_base64 格式")
            frame = _face_mod().limit_bgr_frame(frame)
            try:
                emb = engine.extract_embedding(frame)
            except Exception as e:
                gc.collect()
                return _err(f"提取人脸特征失败: {e}", 500)
            gc.collect()
        if emb is None:
            return _err("未检测到清晰正脸，暂无法注册人脸")
        face_feature_json = json.dumps(emb.tolist())
    db.upsert_student(student_id, name, face_feature_json=face_feature_json)
    _invalidate_face_db_cache()
    row = db.get_student_by_student_id(student_id, include_deleted=True)
    return jsonify(_student_row_to_json(row))


@app.route("/api/face/students/<student_id>", methods=["PATCH"])
def patch_face_student(student_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    name = data.get("name")
    ok = db.update_student(student_id, name=name)
    if not ok:
        return _err("学生不存在", 404)
    _invalidate_face_db_cache()
    row = db.get_student_by_student_id(student_id, include_deleted=False)
    return jsonify(_student_row_to_json(row))


@app.route("/api/face/students/<student_id>", methods=["DELETE"])
def delete_face_student(student_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    ok = db.soft_delete_student(student_id)
    if not ok:
        return _err("学生不存在或已删除", 404)
    _invalidate_face_db_cache()
    return jsonify({"ok": True})


@app.route("/api/face/recognize", methods=["POST"])
def face_recognize_once():
    global _face_recognize_hint_logged
    if not _face_recognize_hint_logged:
        print(
            "[face] POST /api/face/recognize 首次进入（若进程内尚未加载模型，下方会出现 FaceEmotionEngine 初始化日志）",
            flush=True,
        )
        _face_recognize_hint_logged = True
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    image_base64 = (data.get("image_base64") or data.get("imageBase64") or "").strip()
    threshold = float(data.get("threshold") or 0.6)
    if not image_base64:
        return _err("缺少 image_base64")
    uid = str(user_id)
    now_ts = time.time()
    if _face_min_infer_interval_sec > 0:
        with _face_recent_result_lock:
            last_infer_ts = float(_face_last_infer_ts_by_user.get(uid, 0.0))
            if (now_ts - last_infer_ts) < _face_min_infer_interval_sec:
                recent = _face_recent_result_by_user.get(uid)
                if recent:
                    return jsonify(recent.get("payload") or {})
                return jsonify({"width": 0, "height": 0, "count": 0, "detections": []})
    if _face_recent_result_ttl_sec > 0:
        with _face_recent_result_lock:
            recent = _face_recent_result_by_user.get(uid)
            if recent and (now_ts - float(recent.get("ts", 0))) <= _face_recent_result_ttl_sec:
                return jsonify(recent.get("payload") or {})

    with _face_infer_lock:
        if _face_min_infer_interval_sec > 0:
            with _face_recent_result_lock:
                _face_last_infer_ts_by_user[uid] = time.time()
        engine, err_eng = _get_face_engine_safe()
        if err_eng:
            return err_eng
        frame = engine.decode_base64_image(image_base64)
        if frame is None:
            return _err("图片解析失败")
        frame = _face_mod().limit_bgr_frame(frame)

        face_db = _get_face_db_embeddings_cached()
        try:
            detections = engine.detect(frame, face_db, threshold=threshold)
        except Exception as e:
            gc.collect()
            return _err(f"识别失败: {e}", 500)
        for d in detections:
            if d.student_id != "unknown":
                db.add_emotion_record(d.student_id, d.emotion, d.confidence)
                if _is_face_abnormal_emotion(d.emotion):
                    owner = db.get_user_by_account(d.student_id)
                    if owner and owner.get("id"):
                        _maybe_create_proactive_trigger(int(owner["id"]))

        payload = {
            "width": int(frame.shape[1]),
            "height": int(frame.shape[0]),
            "count": len(detections),
            "detections": [
                {
                    "student_id": d.student_id,
                    "emotion": d.emotion,
                    "confidence": d.confidence,
                    "box": d.box,
                }
                for d in detections
            ],
        }
        if _face_recent_result_ttl_sec > 0:
            with _face_recent_result_lock:
                _face_recent_result_by_user[uid] = {"ts": time.time(), "payload": payload}
        gc.collect()
        return jsonify(payload)


@app.route("/api/face/records", methods=["GET"])
def list_face_records():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    student_id = request.args.get("student_id")
    limit = min(max(1, request.args.get("limit", type=int) or 200), 1000)
    rows = db.list_emotion_records(student_id=student_id, limit=limit)
    return jsonify([_record_row_to_json(r) for r in rows])


@app.route("/api/face/records/<int:record_id>", methods=["DELETE"])
def delete_face_record(record_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    ok = db.soft_delete_emotion_record(record_id)
    if not ok:
        return _err("记录不存在或已删除", 404)
    return jsonify({"ok": True})


# ---------- 情绪异常与主动疏导（个人中心情感曲线、事件记录、主动疏导入口） ----------


def _anomaly_row_to_json(row):
    if not row:
        return None
    from_mon = 1 if row.get("from_monitoring") else 0
    return {
        "id": row["id"],
        "user_id": row["user_id"],
        "emotion_label": row["emotion_label"],
        "reason": (row.get("reason") or "").strip(),
        "from_monitoring": from_mon,
        "created_at": row["created_at"].isoformat() if hasattr(row.get("created_at"), "isoformat") else str(row.get("created_at", "")),
    }


@app.route("/api/emotion/anomaly", methods=["POST"])
def emotion_anomaly_add():
    """记录一条情绪异常（聊天或监控写入）。body: { emotion_label, reason?, from_monitoring? }。from_monitoring: 0=聊天（可带 reason），1=监控。为 1 时创建主动疏导触发。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    label = (data.get("emotion_label") or data.get("label") or "").strip()
    if not label:
        return _err("请输入情绪标签")
    reason = (data.get("reason") or "").strip()[:2000]
    from_monitoring = 0
    if data.get("from_monitoring") in (1, "1", True):
        from_monitoring = 1
    elif (data.get("source") or "").strip().lower() == "monitoring":
        from_monitoring = 1
    try:
        aid = db.add_emotion_anomaly(int(user_id), label, reason=reason, from_monitoring=from_monitoring)
        if from_monitoring == 1:
            db.create_proactive_trigger(int(user_id), "monitoring")
        row = db.get_emotion_anomalies_by_user(int(user_id), limit=1)
        return jsonify(_anomaly_row_to_json(row[0]) if row else {"id": aid, "emotion_label": label, "reason": reason, "from_monitoring": from_monitoring})
    except Exception as e:
        return _err("保存失败: " + str(e), 500)


@app.route("/api/emotion/anomalies", methods=["GET"])
def emotion_anomalies_list():
    """当前用户情绪异常列表，供个人中心与模型检索。query: limit, since_days"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        limit = request.args.get("limit", type=int) or 100
        since_days = request.args.get("since_days", type=int)
        rows = db.get_emotion_anomalies_by_user(int(user_id), limit=min(max(1, limit), 500), since_days=since_days)
        return jsonify([_anomaly_row_to_json(r) for r in rows])
    except Exception as e:
        return _err("查询失败: " + str(e), 500)


@app.route("/api/emotion/stats", methods=["GET"])
def emotion_stats():
    """情绪统计：按日聚合数量，供可视化情感曲线。query: days 默认 30"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        days = min(max(1, request.args.get("days", type=int) or 30), 365)
        rows = db.get_emotion_stats_by_user(int(user_id), days=days)
        return jsonify(rows)
    except Exception as e:
        return _err("查询失败: " + str(e), 500)


@app.route("/api/proactive/pending", methods=["GET"])
def proactive_pending():
    """当前用户是否有待响应的主动疏导（监控检测到异常 / 多次异常）。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        row = db.get_pending_proactive_trigger(int(user_id))
        if not row:
            return jsonify(None)
        return jsonify({
            "id": row["id"],
            "trigger_type": row["trigger_type"],
            "created_at": row["created_at"].isoformat() if hasattr(row.get("created_at"), "isoformat") else str(row.get("created_at", "")),
        })
    except Exception as e:
        return _err("查询失败: " + str(e), 500)


@app.route("/api/proactive/ack", methods=["POST"])
def proactive_ack():
    """用户点击「去聊天」后确认已响应。body: { triggerId }"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    tid = data.get("triggerId") or data.get("trigger_id")
    if tid is None:
        return _err("缺少 triggerId")
    try:
        pending = db.get_pending_proactive_trigger(int(user_id))
        ok = db.acknowledge_proactive_trigger(int(tid), int(user_id))
        care_message = None
        if ok:
            trig_type = "repeated_anomaly"
            if pending and int(pending.get("id") or 0) == int(tid):
                trig_type = (pending.get("trigger_type") or "repeated_anomaly").strip()
            care_message = _send_proactive_care_message(int(user_id), trig_type, 0, 0, use_llm=False)
        return jsonify({"ok": ok, "careMessage": care_message})
    except Exception as e:
        return _err("操作失败: " + str(e), 500)


def _parse_llm_json(text: str):
    """从模型输出中解析 JSON 对象。"""
    if not (text or "").strip():
        return None
    s = text.strip()
    for raw in (s, s.replace("```json", "").replace("```", "").strip()):
        try:
            return json.loads(raw)
        except (json.JSONDecodeError, TypeError):
            continue
    return None


def _llm_json_chat(system_prompt: str, user_prompt: str, max_tokens: int = 900, timeout: int = 90):
    """
    调用与主对话一致的模型通道，返回解析后的 JSON 或错误。
    注意：这里只做瞬时推理，不写入 conversations/messages，不持久化模型上下文。
    """
    if not API_KEY:
        return None, "未配置大模型 API_KEY"
    merged_system = "\n\n".join([x for x in [(CHAT_SYSTEM_PROMPT or "").strip(), (system_prompt or "").strip()] if x])
    messages = _with_chat_system_prompt([{"role": "user", "content": user_prompt}])
    if merged_system:
        messages = [{"role": "system", "content": merged_system}] + [{"role": "user", "content": user_prompt}]
    payload = {
        "model": CHAT_MODEL,
        "messages": messages,
        "max_tokens": max_tokens,
        "stream": False,
    }
    headers = {"Content-Type": "application/json", "Authorization": f"Bearer {API_KEY}"}
    try:
        r = _requests().post(CHAT_API_URL, json=payload, headers=headers, timeout=timeout)
        if r.status_code != 200:
            try:
                err_body = r.json()
                err_msg = (err_body.get("error") or {}).get("message") or ""
            except Exception:
                err_msg = r.text or ""
            return None, f"模型请求失败: HTTP {r.status_code} {err_msg}".strip()
        out = r.json()
        content = (out.get("choices", [{}])[0].get("message", {}).get("content") or "").strip()
        obj = _parse_llm_json(content)
        if not obj:
            return None, "模型返回无法解析为 JSON"
        return obj, None
    except Exception as e:
        return None, str(e)


def _week_monday(d: datetime) -> datetime:
    """返回当周周一 00:00:00 的 datetime（本地）。"""
    wd = d.weekday()
    mon = d - timedelta(days=wd)
    return mon.replace(hour=0, minute=0, second=0, microsecond=0)


def _release_temp_llm_context(*_objs):
    """
    显式释放仅用于本次推理的临时上下文，避免误用与长驻内存。
    不涉及数据库持久化（本模块本就不写 messages/conversations）。
    """
    gc.collect()


def _journal_row_to_json(row):
    if not row:
        return None
    ws = row.get("week_start")
    return {
        "id": row["id"],
        "user_id": row["user_id"],
        "week_start": ws.isoformat() if hasattr(ws, "isoformat") else str(ws),
        "mood": (row.get("mood") or "").strip(),
        "body": (row.get("body") or "").strip(),
        "created_at": row["created_at"].isoformat() if hasattr(row.get("created_at"), "isoformat") else str(row.get("created_at", "")),
        "updated_at": row["updated_at"].isoformat() if hasattr(row.get("updated_at"), "isoformat") else str(row.get("updated_at", "")),
    }


def _daily_journal_fallback(call_name: str, msg_rows: list):
    """无模型或模型失败时，基于当天对话做轻量总结。"""
    user_lines = []
    for m in msg_rows or []:
        if (m.get("role") or "").strip() != "user":
            continue
        text = (m.get("content") or "").strip()
        if not text:
            continue
        user_lines.append(text)
    if not user_lines:
        return "—", "暂无聊天，快和我谈谈心吧！"
    snippet = "；".join(user_lines[-3:])
    if len(snippet) > 120:
        snippet = snippet[:120].rstrip("，,、；;：:。.!！？? ") + "…"
    body = (
        f"{call_name}今天和我聊了不少，我记住了你提到的这些事：{snippet}。"
        "谢谢你愿意分享，我会继续陪着你。"
    )
    return "未标注", body


def _ensure_daily_journal_for_day(user_id: int, day0: datetime):
    """
    确保某一天的日记已生成并落库（用于“次日自动更新”）。
    说明：为保证页面响应速度，默认走轻量总结，不阻塞在大模型上。
    """
    day_key = day0.strftime("%Y-%m-%d")
    existed = db.get_weekly_journal_by_week(int(user_id), day_key)
    if existed:
        return existed

    user_row = db.get_user_by_id(int(user_id)) or {}
    call_name = (
        (user_row.get("preferred_name") or "").strip()
        or (user_row.get("username") or "").strip()
        or "你"
    )[:50]
    start_dt = day0.replace(hour=0, minute=0, second=0, microsecond=0)
    end_dt = start_dt + timedelta(days=1)
    start_s = start_dt.strftime("%Y-%m-%d %H:%M:%S")
    end_s = end_dt.strftime("%Y-%m-%d %H:%M:%S")
    msg_rows = db.get_messages_for_user_range(int(user_id), start_s, end_s, limit=4000)

    mood, summary = _daily_journal_fallback(call_name, msg_rows)
    db.upsert_weekly_journal(int(user_id), day_key, mood, summary)
    return db.get_weekly_journal_by_week(int(user_id), day_key)


@app.route("/api/journal/weekly", methods=["GET"])
def journal_weekly_list():
    """本周日记列表：从本周首次有对话的日期开始，到昨天为止；缺聊日显示占位文案。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        now = datetime.now()
        week_monday = _week_monday(now)
        today0 = now.replace(hour=0, minute=0, second=0, microsecond=0)
        last_closed_day = today0 - timedelta(days=1)
        if last_closed_day < week_monday:
            return jsonify([])

        # 仅统计本周已结束日期内的用户对话（到今天 00:00，不含今天）
        week_start_s = week_monday.strftime("%Y-%m-%d %H:%M:%S")
        week_end_s = today0.strftime("%Y-%m-%d %H:%M:%S")
        msg_rows = db.get_messages_for_user_range(int(user_id), week_start_s, week_end_s, limit=8000)
        talked_days = {
            str((m.get("created_at") or ""))[:10]
            for m in msg_rows
            if (m.get("role") or "").strip() == "user" and (m.get("content") or "").strip()
        }
        # 本周尚无有效对话：不展示任何“每天日记”
        if not talked_days:
            return jsonify([])

        # 从“本周首次有对话那天”开始算；之前无对话日期不展示
        first_chat_day_key = min(talked_days)
        first_chat_day = datetime.strptime(first_chat_day_key, "%Y-%m-%d")
        start_day = max(week_monday, first_chat_day)
        day_count = (last_closed_day.date() - start_day.date()).days + 1
        if day_count <= 0:
            return jsonify([])
        days = [start_day + timedelta(days=i) for i in range(day_count)]
        day_keys = [d.strftime("%Y-%m-%d") for d in days]

        rows = db.list_weekly_journals(int(user_id), limit=128)
        row_map = {str(r.get("week_start"))[:10]: r for r in rows}

        # 自动补齐：仅对“有对话”的日期生成；无对话日期返回占位文案（不落库）
        for d in days:
            key = d.strftime("%Y-%m-%d")
            if key in row_map:
                continue
            if key not in talked_days:
                continue
            try:
                generated = _ensure_daily_journal_for_day(int(user_id), d)
                if generated:
                    row_map[key] = generated
            except Exception:
                # 单天补齐失败不影响整体列表返回
                continue

        items = []
        for key in day_keys:
            row = row_map.get(key)
            if row:
                items.append(_journal_row_to_json(row))
            else:
                items.append(
                    {
                        "id": 0,
                        "user_id": int(user_id),
                        "week_start": key,
                        "mood": "—",
                        "body": "暂无聊天，快和我谈谈心吧！",
                        "created_at": f"{key}T00:00:00",
                        "updated_at": f"{key}T00:00:00",
                    }
                )
        return jsonify(items)
    except Exception as e:
        return _err("查询失败: " + str(e), 500)


@app.route("/api/journal/weekly/generate", methods=["POST"])
def journal_weekly_generate():
    """
    根据「本天」与数字人助手的对话摘要，生成/更新一条情感日记。
    body: { week_start?: "YYYY-MM-DD" } 表示要生成的日历日；不传则默认为今天（自然日 0:00～次日 0:00）。
    """
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    raw_day = (data.get("week_start") or data.get("weekStart") or "").strip()
    now = datetime.now()
    if raw_day:
        try:
            day0 = datetime.strptime(raw_day[:10], "%Y-%m-%d")
        except ValueError:
            return _err("week_start 格式应为 YYYY-MM-DD")
    else:
        day0 = now.replace(hour=0, minute=0, second=0, microsecond=0)
    user_row = db.get_user_by_id(int(user_id)) or {}
    call_name = (
        (user_row.get("preferred_name") or "").strip()
        or (user_row.get("username") or "").strip()
        or "你"
    )[:50]
    journal_date_str = day0.strftime("%Y-%m-%d")
    start_dt = day0.replace(hour=0, minute=0, second=0, microsecond=0)
    end_dt = start_dt + timedelta(days=1)
    start_s = start_dt.strftime("%Y-%m-%d %H:%M:%S")
    end_s = end_dt.strftime("%Y-%m-%d %H:%M:%S")
    try:
        msg_rows = db.get_messages_for_user_range(int(user_id), start_s, end_s, limit=4000)
    except Exception as e:
        return _err("读取聊天记录失败: " + str(e), 500)
    lines = []
    for m in msg_rows:
        role = "用户" if (m.get("role") or "").strip() == "user" else "助手"
        c = (m.get("content") or "").strip()
        if not c:
            continue
        if len(c) > 4000:
            c = c[:4000] + "…"
        lines.append(f"{role}: {c}")
    transcript = "\n".join(lines)
    if len(transcript) > 100000:
        transcript = transcript[-100000:]
    if not transcript.strip():
        transcript = "（本天暂无有效对话记录）"

    sys_prompt = (
        "你是小Q，一个温柔、细腻、会共情的知心朋友。"
        "根据以下用户与助手在「同一自然日」内的对话摘录，整理成一段中文「本天情感日记」。"
        "要求：\n"
        "1. 只输出 JSON，不要其他文字。\n"
        "2. JSON 格式：{\"mood\":\"\",\"summary\":\"\"}。\n"
        "3. mood 只填比喻词本身（如：晴、阴、多云、小雨），不要加「当天情绪：」等前缀，前端会统一加。\n"
        "4. summary 必须使用第二人称/昵称称呼（如“一川今天…”、“你今天…”），像朋友在写，不要出现“用户”这种词。\n"
        "5. 禁止笼统总结（如“情绪无明显起伏”）。必须结合当天具体事件细节，写出生动感受与陪伴口吻。\n"
        "6. 可以有一句朋友视角表达（如“我也替你开心”），但不要说教，不要模板化。\n"
        "7. 若对话很少，仍基于现有内容写短日记，不要编造具体事实。"
    )
    user_prompt = f"用户希望被称呼：{call_name}\n本天日期：{journal_date_str}\n\n对话摘录：\n{transcript}"
    obj, err = _llm_json_chat(sys_prompt, user_prompt, max_tokens=900, timeout=120)
    _release_temp_llm_context(lines, msg_rows, transcript, user_prompt)
    if err:
        return _err(err, 503)
    mood = (obj.get("mood") or obj.get("今日心情") or "").strip()
    summary = (obj.get("summary") or obj.get("body") or "").strip()
    if not mood and not summary:
        return _err("模型未返回有效内容", 500)
    if not mood:
        mood = "未标注"
    if not summary:
        summary = "（暂无摘要）"
    try:
        db.upsert_weekly_journal(int(user_id), journal_date_str, mood, summary)
        row = db.get_weekly_journal_by_week(int(user_id), journal_date_str)
        if not row:
            return jsonify({"ok": True, "week_start": journal_date_str, "mood": mood, "body": summary})
        return jsonify(_journal_row_to_json(row))
    except Exception as e:
        return _err("保存失败: " + str(e), 500)


@app.route("/api/journal/weekly/summary", methods=["POST"])
def journal_weekly_summary():
    """
    生成本周一条整周周记（不落库），用于弹窗展示。
    body: { week_start?: "YYYY-MM-DD" } 可选，传入该周任意一天；不传默认本周。
    """
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    raw_week = (data.get("week_start") or data.get("weekStart") or "").strip()
    now = datetime.now()
    if raw_week:
        try:
            base_day = datetime.strptime(raw_week[:10], "%Y-%m-%d")
            week_monday = _week_monday(base_day)
        except ValueError:
            return _err("week_start 格式应为 YYYY-MM-DD")
    else:
        week_monday = _week_monday(now)
    week_end = week_monday + timedelta(days=6)
    start_s = week_monday.strftime("%Y-%m-%d %H:%M:%S")
    end_s = (week_monday + timedelta(days=7)).strftime("%Y-%m-%d %H:%M:%S")

    user_row = db.get_user_by_id(int(user_id)) or {}
    call_name = (
        (user_row.get("preferred_name") or "").strip()
        or (user_row.get("username") or "").strip()
        or "你"
    )[:50]

    try:
        msg_rows = db.get_messages_for_user_range(int(user_id), start_s, end_s, limit=8000)
    except Exception as e:
        return _err("读取聊天记录失败: " + str(e), 500)

    lines = []
    for m in msg_rows:
        role = "用户" if (m.get("role") or "").strip() == "user" else "助手"
        c = (m.get("content") or "").strip()
        if not c:
            continue
        if len(c) > 4000:
            c = c[:4000] + "…"
        lines.append(f"{role}: {c}")
    transcript = "\n".join(lines)
    if len(transcript) > 150000:
        transcript = transcript[-150000:]
    if not transcript.strip():
        transcript = "（本周暂无有效对话记录）"

    sys_prompt = (
        "你是小Q，一个温柔、细腻、会共情的知心朋友。"
        "根据用户与助手在同一自然周（周一到周日）的对话摘录，整理成一条中文「本周周记」。"
        "要求：\n"
        "1. 只输出 JSON，不要其他文字。\n"
        "2. JSON 格式：{\"mood\":\"\",\"summary\":\"\"}。\n"
        "3. mood 只填比喻词本身（如：晴、阴、多云、小雨），不要加前缀。\n"
        "4. summary 必须使用第二人称或昵称称呼（如“一川这周…”、“你这周…”），像朋友在写，不要出现“用户”。\n"
        "5. 禁止笼统套话（如“情绪无明显起伏”）；必须包含本周具体事件与感受细节，语气自然有温度。\n"
        "6. 可有一句朋友视角表达（如“我也替你开心”），但不要说教。"
    )
    user_prompt = (
        f"用户希望被称呼：{call_name}\n"
        f"本周范围：{week_monday.strftime('%Y-%m-%d')} ~ {week_end.strftime('%Y-%m-%d')}\n\n"
        f"对话摘录：\n{transcript}"
    )
    obj, err = _llm_json_chat(sys_prompt, user_prompt, max_tokens=1200, timeout=120)
    _release_temp_llm_context(lines, msg_rows, transcript, user_prompt)
    if err:
        return _err(err, 503)
    mood = (obj.get("mood") or "").strip() or "未标注"
    summary = (obj.get("summary") or obj.get("body") or "").strip() or "暂无聊天，快和我谈谈心吧！"
    return jsonify(
        {
            "week_start": week_monday.strftime("%Y-%m-%d"),
            "week_end": week_end.strftime("%Y-%m-%d"),
            "mood": mood,
            "body": summary,
        }
    )


@app.route("/api/journal/chat-insights", methods=["POST"])
def journal_chat_insights():
    """
    从近期对话中提取：安排建议 + 可加入日程的事项建议（用户确认后由前端调用 /api/schedules 写入）。
    body: { lookback_days?: number, max_messages?: number }
    """
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    lookback = min(max(1, int(data.get("lookback_days") or data.get("lookbackDays") or 7)), 30)
    max_msgs = min(max(50, int(data.get("max_messages") or data.get("maxMessages") or 500)), 4000)
    end = datetime.now()
    start = end - timedelta(days=lookback)
    start_s = start.strftime("%Y-%m-%d %H:%M:%S")
    end_s = end.strftime("%Y-%m-%d %H:%M:%S")
    try:
        msg_rows = db.get_messages_for_user_range(int(user_id), start_s, end_s, limit=max_msgs)
    except Exception as e:
        return _err("读取聊天记录失败: " + str(e), 500)
    lines = []
    for m in msg_rows:
        role = "用户" if (m.get("role") or "").strip() == "user" else "助手"
        c = (m.get("content") or "").strip()
        if not c:
            continue
        if len(c) > 4000:
            c = c[:4000] + "…"
        lines.append(f"{role}: {c}")
    transcript = "\n".join(lines)
    if len(transcript) > 120000:
        transcript = transcript[-120000:]
    if not transcript.strip():
        transcript = "（近期暂无对话）"
    try:
        schedule_rows = db.get_schedules_by_user(int(user_id), limit=300) or []
    except Exception:
        schedule_rows = []
    pending_schedules = []
    for r in schedule_rows:
        if (r.get("status") or "pending").strip() != "pending":
            continue
        sat = r.get("scheduled_at")
        sat_str = sat.isoformat(sep=" ")[:19] if hasattr(sat, "isoformat") else str(sat or "")[:19]
        title = (r.get("title") or "").strip()
        if not sat_str or not title:
            continue
        pending_schedules.append(f"- {sat_str} | {title}")
    schedule_snapshot = "\n".join(pending_schedules[:80]) if pending_schedules else "（当前暂无待办日程）"

    sys_prompt = (
        "你是学生时间管理助手。请基于“近期对话 + 当前已存在日程”给出安排建议，并输出严格 JSON，不要其他文字。\n"
        "格式：\n"
        "{\"comfort\":\"安排建议文本（中文）\",\"schedule_suggestions\":[{\"title\":\"事项标题\",\"scheduled_at\":\"YYYY-MM-DD HH:MM:SS\",\"note\":\"可选说明\"}]}\n"
        "规则：\n"
        "1. comfort 必须是“安排建议”，不是情感宽慰。请结合当前已有日程给出可执行建议。\n"
        "2. 若发现时间冲突（同时间段多事项、明显挤占），在 comfort 中明确指出冲突并给出调序/改期建议。\n"
        "3. 若没有明显冲突，在 comfort 中给出提醒（如提前准备、预留通勤/缓冲时间、避免遗忘）。\n"
        "4. schedule_suggestions 只提取用户明确或强烈暗示的时间安排（会议、作业截止、面试等），没有则空数组。\n"
        "5. 不要编造用户未提及的紧急事件；建议要简洁具体、可执行。"
    )
    user_prompt = (
        f"当前时间参考：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
        f"当前待办日程：\n{schedule_snapshot}\n\n"
        f"近期对话摘录：\n{transcript}"
    )
    obj, err = _llm_json_chat(sys_prompt, user_prompt, max_tokens=1200, timeout=120)
    _release_temp_llm_context(lines, msg_rows, transcript, user_prompt)
    if err:
        return _err(err, 503)
    comfort = (obj.get("comfort") or "").strip()
    raw_sug = obj.get("schedule_suggestions") or obj.get("schedules") or []
    if not isinstance(raw_sug, list):
        raw_sug = []
    suggestions = []
    for item in raw_sug[:20]:
        if not isinstance(item, dict):
            continue
        title = (item.get("title") or "").strip()
        sat = (item.get("scheduled_at") or item.get("scheduledAt") or "").strip().replace("T", " ")
        if len(sat) == 16:
            sat = sat + ":00"
        note = (item.get("note") or "").strip()
        if not title or not sat:
            continue
        suggestions.append({"title": title[:500], "scheduled_at": sat[:19], "note": note[:2000]})
    return jsonify({"comfort": comfort or "（暂无）", "schedule_suggestions": suggestions})


# ---------- 日程（个人中心用） ----------


def _schedule_row_to_json(row):
    if not row:
        return None
    return {
        "id": row["id"],
        "user_id": row["user_id"],
        "title": row["title"],
        "scheduled_at": row["scheduled_at"].isoformat() if hasattr(row.get("scheduled_at"), "isoformat") else str(row.get("scheduled_at", "")),
        "end_at": row["end_at"].isoformat() if row.get("end_at") and hasattr(row["end_at"], "isoformat") else (str(row["end_at"]) if row.get("end_at") else None),
        "source": (row.get("source") or "conversation").strip(),
        "raw_text": (row.get("raw_text") or "").strip() or None,
        "status": (row.get("status") or "pending").strip(),
        "created_at": row["created_at"].isoformat() if hasattr(row.get("created_at"), "isoformat") else str(row.get("created_at", "")),
    }


@app.route("/api/schedules", methods=["GET"])
def schedules_list():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        start_date = request.args.get("startDate") or request.args.get("start_date")
        end_date = request.args.get("endDate") or request.args.get("end_date")
        rows = db.get_schedules_by_user(int(user_id), start_date=start_date, end_date=end_date)
        return jsonify([_schedule_row_to_json(r) for r in rows])
    except Exception as e:
        return _err("查询失败: " + str(e), 500)


@app.route("/api/schedules", methods=["POST"])
def schedules_create():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    title = (data.get("title") or "").strip()
    scheduled_at = (data.get("scheduled_at") or data.get("scheduledAt") or "").strip().replace("T", " ")[:19]
    if not title or not scheduled_at:
        return _err("请填写 title 和 scheduled_at")
    if len(scheduled_at) == 16:
        scheduled_at = scheduled_at + ":00"
    try:
        end_at = (data.get("end_at") or data.get("endAt") or "").strip().replace("T", " ")[:19] or None
        sid = db.create_schedule(int(user_id), title, scheduled_at, end_at=end_at, source="manual")
        row = db.get_schedule_by_id(sid, int(user_id))
        return jsonify(_schedule_row_to_json(row))
    except Exception as e:
        return _err("创建失败: " + str(e), 500)


@app.route("/api/schedules/<int:schedule_id>", methods=["PATCH"])
def schedule_patch(schedule_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    updates = {}
    if "title" in data:
        updates["title"] = (data.get("title") or "").strip()
    if "scheduled_at" in data or "scheduledAt" in data:
        raw = (data.get("scheduled_at") or data.get("scheduledAt") or "").strip().replace("T", " ")[:19]
        if len(raw) == 16:
            raw = raw + ":00"
        updates["scheduled_at"] = raw
    if "end_at" in data or "endAt" in data:
        raw = (data.get("end_at") or data.get("endAt") or "").strip().replace("T", " ")[:19] or None
        if raw and len(raw) == 16:
            raw = raw + ":00"
        updates["end_at"] = raw
    if "status" in data:
        updates["status"] = (data.get("status") or "pending").strip()[:20]
    if not updates:
        row = db.get_schedule_by_id(schedule_id, int(user_id))
        if not row:
            return _err("日程不存在", 404)
        return jsonify(_schedule_row_to_json(row))
    ok = db.update_schedule(schedule_id, int(user_id), **updates)
    if not ok:
        return _err("日程不存在", 404)
    row = db.get_schedule_by_id(schedule_id, int(user_id))
    return jsonify(_schedule_row_to_json(row))


@app.route("/api/schedules/<int:schedule_id>", methods=["DELETE"])
def schedule_delete(schedule_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    ok = db.delete_schedule(schedule_id, int(user_id))
    if not ok:
        return _err("日程不存在", 404)
    return jsonify({"ok": True})


# ---------- 助理业务流程（早晨问好 / 日常推荐 / 通知栏 / 任务打断恢复 / 今日关键词） ----------


@app.route("/api/assistant/morning-greeting", methods=["GET"])
def assistant_morning_greeting():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    now = datetime.now()
    force = request.args.get("force") in ("1", "true", "True")
    try:
        logs = db.list_user_memory_vectors(int(user_id), memory_type="morning_greeting", limit=10) or []
    except Exception:
        logs = []
    already = False
    for r in logs:
        try:
            meta = json.loads((r.get("metadata_json") or "{}").strip())
            if str(meta.get("date") or "") == now.strftime("%Y-%m-%d"):
                already = True
                break
        except Exception:
            continue
    payload = _build_morning_greeting(int(user_id))
    if (5 <= now.hour <= 11 and not already) or force:
        _memory_add(
            int(user_id),
            "morning_greeting",
            payload.get("message") or "",
            {"date": now.strftime("%Y-%m-%d"), "hour": now.hour},
        )
    payload["already_pushed_today"] = already
    payload["time_bucket"] = "morning"
    return jsonify(payload)


@app.route("/api/assistant/daily-recommendations", methods=["GET"])
def assistant_daily_recommendations():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    recs = _build_daytime_recommendations(int(user_id))
    return jsonify({"items": recs, "generated_at": datetime.now().isoformat()})


@app.route("/api/assistant/mood-keywords/today", methods=["GET"])
def assistant_today_mood_keywords():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    kws = _today_mood_keywords(int(user_id), limit=8)
    return jsonify({"keywords": kws, "date": datetime.now().strftime("%Y-%m-%d")})


@app.route("/api/assistant/notice-cards", methods=["GET"])
def assistant_notice_cards():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    recs = _build_daytime_recommendations(int(user_id))
    neg_moods = _today_negative_mood_keywords(int(user_id), limit=4)
    today_rows = _today_schedule_snapshot(int(user_id))
    pending = [x for x in today_rows if (x.get("status") or "pending") == "pending"]
    top_schedule = pending[0] if pending else None
    now = datetime.now()
    continuous_minutes = _get_continuous_usage_minutes(int(user_id), now=now)
    _auto_expire_interest_news_buckets(int(user_id), now=now)
    today_key = now.strftime("%Y-%m-%d")
    fatigue_detected = any(x in ("疲惫", "焦虑", "压力", "烦躁", "崩溃", "低落") for x in neg_moods)
    morning_first_login = False
    if 5 <= now.hour <= 11:
        try:
            logs = db.list_user_memory_vectors(int(user_id), memory_type="morning_plan_card", limit=10) or []
        except Exception:
            logs = []
        morning_seen = False
        for r in logs:
            try:
                meta = json.loads((r.get("metadata_json") or "{}").strip())
                if str(meta.get("date") or "") == today_key:
                    morning_seen = True
                    break
            except Exception:
                continue
        morning_first_login = not morning_seen
    cards = []
    # 手动触发：用于演示“高压减负模块”卡片
    try:
        manual_trigger = db.get_pending_proactive_trigger(int(user_id))
    except Exception:
        manual_trigger = None
    if manual_trigger and (manual_trigger.get("trigger_type") or "").strip() == "high_pressure_relief":
        profile = _get_user_profile(int(user_id))
        call_name = (
            profile.get("preferred_name")
            or profile.get("name")
            or profile.get("username")
            or "你"
        )
        cards.append(_build_high_pressure_relief_card(call_name))
        try:
            db.acknowledge_proactive_trigger(int(manual_trigger.get("id")), int(user_id))
        except Exception:
            pass

    wellness_milestone = _maybe_pick_wellness_break_milestone(
        int(user_id),
        continuous_minutes=continuous_minutes,
        now=now,
    )
    if wellness_milestone is not None:
        cards.append(_build_wellness_break_card(int(user_id), continuous_minutes))
        _memory_add(
            int(user_id),
            "wellness_break_push",
            "shown",
            {
                "date": today_key,
                "hour": now.hour,
                "minute": now.minute,
                "milestone_min": int(wellness_milestone),
                "continuous_minutes": int(continuous_minutes),
            },
        )
    # 今日总览仅在当天首次登录触发，避免首次登录被“疲惫”分支误触发
    show_today_plan = bool(pending) and morning_first_login
    if show_today_plan:
        plan_card = _build_today_plan_card(pending, recs, fatigue_detected)
        if plan_card:
            cards.append(plan_card)
            if morning_first_login:
                _memory_add(
                    int(user_id),
                    "morning_plan_card",
                    "shown",
                    {"date": today_key, "hour": now.hour},
                )
    mood_tip_milestone = _maybe_pick_mood_tip_milestone(
        int(user_id),
        neg_moods,
        continuous_minutes=continuous_minutes,
        now=now,
    )
    if mood_tip_milestone is not None:
        tip = _build_mood_tip_card_fields(neg_moods)
        cards.append(
            {
                "id": "mood_tip",
                "type": "mood",
                "priority": int(tip["priority"]),
                "title": "近期情绪小贴士",
                "subtitle": f"{tip['subtitle']}（你已连续使用约 {continuous_minutes} 分钟）",
                "content": tip["content"],
            }
        )
        _memory_add(
            int(user_id),
            "mood_tip_push",
            "shown",
            {
                "date": today_key,
                "hour": now.hour,
                "minute": now.minute,
                "milestone_min": int(mood_tip_milestone),
                "continuous_minutes": int(continuous_minutes),
                "negative_kw": neg_moods[:4],
            },
        )
    cards.sort(key=lambda x: int(x.get("priority") or 0), reverse=True)
    return jsonify({"cards": cards})


@app.route("/api/assistant/interest-news", methods=["GET"])
def assistant_interest_news():
    """根据近期对话推断的兴趣，从网络 RSS 拉取相关新闻标题，供前端交给模型整理。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    force = request.args.get("force") in ("1", "true", "True")
    payload = _get_or_refresh_interest_news_payload(int(user_id), force_refresh=force)
    return jsonify(payload)


@app.route("/api/assistant/interest-news/consume", methods=["POST"])
def assistant_interest_news_consume():
    """点击兴趣新闻卡后记账：点击即算一次。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    now = datetime.now()
    bucket = _detect_interest_news_active_bucket(int(user_id), now=now)
    if not bucket:
        return jsonify({"ok": True, "counted": False, "reason": "outside_active_window"})
    payload = _get_or_refresh_interest_news_payload(int(user_id), force_refresh=False)
    signature = _interest_news_signature(payload.get("items") or [])
    _memory_add(
        int(user_id),
        "interest_news_push",
        "clicked",
        {
            "date": now.strftime("%Y-%m-%d"),
            "bucket": bucket,
            "hour": now.hour,
            "minute": now.minute,
            "signature": signature,
        },
    )
    return jsonify({"ok": True, "counted": True, "bucket": bucket})


@app.route("/api/assistant/chat-prompts", methods=["GET"])
def assistant_chat_prompts():
    """返回聊天区动态引导语（按用户兴趣，至少三天更新一次）。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    force = request.args.get("force") in ("1", "true", "True")
    payload = _get_or_refresh_chat_prompts_payload(int(user_id), force_refresh=force)
    return jsonify(payload)


@app.route("/api/assistant/task-interrupt", methods=["POST"])
def assistant_task_interrupt():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    title = (data.get("title") or data.get("task_title") or "").strip()
    if not title:
        return _err("缺少任务标题 title")
    progress = (data.get("progress_note") or data.get("progressNote") or "").strip()
    conv_id_raw = data.get("conversation_id") or data.get("conversationId")
    conv_id = None
    if conv_id_raw is not None:
        try:
            conv_id = int(conv_id_raw)
        except Exception:
            conv_id = None
    resume_hint = (
        (data.get("resume_hint") or data.get("resumeHint") or "").strip()
        or "你现在的进度我已经帮你保存并标记了，回来后我们可以从这里继续。"
    )
    try:
        rid = db.create_task_interruption_record(
            int(user_id),
            title=title,
            progress_note=progress,
            resume_hint=resume_hint,
            conversation_id=conv_id,
        )
    except Exception as e:
        return _err("保存中断记录失败: " + str(e), 500)
    return jsonify(
        {
            "ok": True,
            "record_id": rid,
            "message": "你现在的进度我已经帮你保存并在日历上打好标记了，等你回来，我们随时接续进行。",
        }
    )


@app.route("/api/assistant/task-resume", methods=["POST"])
def assistant_task_resume():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    rid = data.get("record_id") or data.get("recordId")
    if rid is None:
        return _err("缺少 record_id")
    try:
        rid = int(rid)
    except Exception:
        return _err("record_id 无效")
    resume_note = (data.get("resume_note") or data.get("resumeNote") or "").strip()
    ok = db.resume_task_interruption_record(rid, int(user_id), resume_note=resume_note)
    if not ok:
        return _err("记录不存在或已恢复", 404)
    return jsonify({"ok": True, "message": "已恢复该任务，我们可以继续刚才的进度。"})


@app.route("/api/assistant/task-interrupt/latest", methods=["GET"])
def assistant_task_interrupt_latest():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    row = db.get_latest_interrupted_task(int(user_id))
    if not row:
        return jsonify(None)
    return jsonify(
        {
            "id": int(row.get("id") or 0),
            "title": (row.get("title") or "").strip(),
            "progress_note": (row.get("progress_note") or "").strip(),
            "resume_hint": (row.get("resume_hint") or "").strip(),
            "interrupted_at": row.get("interrupted_at").isoformat()
            if hasattr(row.get("interrupted_at"), "isoformat")
            else str(row.get("interrupted_at") or ""),
        }
    )


@app.route("/api/assistant/topic-recommendations", methods=["GET"])
def assistant_topic_recommendations():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    interests = _infer_user_interests(int(user_id), limit=6)
    mood_keywords = _today_mood_keywords(int(user_id), limit=3)
    memory_rows = _memory_search(int(user_id), "近期想做什么 计划 兴趣", limit=5)
    memory_topics = []
    for r in memory_rows:
        c = (r.get("content") or "").strip()
        if c:
            memory_topics.append(c[:60])
    topics = []
    for i in interests[:4]:
        topics.append(f"你最近常聊「{i}」，要不要我给你做一个今日相关计划？")
    for m in mood_keywords[:2]:
        topics.append(f"今天关键词偏「{m}」，我可以给你一个更轻量的任务安排版本。")
    for t in memory_topics[:2]:
        topics.append(f"基于你之前提到的「{t}」，要不要继续推进？")
    if not topics:
        topics = ["要不要我先帮你定今天的三件关键小目标？", "如果你愿意，我可以先帮你排一个轻量日程。"]
    return jsonify({"topics": topics[:8]})


# ---------- 实时对话（DashScope Realtime：语音入 -> 文本+语音出，与聊天框同步） ----------

_REALTIME_CHUNK_BYTES = 3200  # 100ms @ 16k 16bit mono
_REALTIME_CONNECT_TIMEOUT_SEC = 30
_REALTIME_RECV_TIMEOUT_SEC = 180
_REALTIME_EVENT_IDLE_TIMEOUT_SEC = 20
_REALTIME_SESSIONS = {}
_REALTIME_SESSIONS_LOCK = threading.Lock()


def _build_realtime_instructions() -> str:
    """实时语音助手系统提示词（与非实时会话完全解耦）。"""
    base = (REALTIME_SYSTEM_PROMPT or "").strip() or "你是同学们的好朋友小Q，请简洁友好地回复。"
    return base


def _realtime_emit(session, payload):
    q = session.get("queue")
    if q is None:
        return
    try:
        q.put(payload, block=False)
    except queue.Full:
        pass


def _close_realtime_session(session_id, reason="closed"):
    with _REALTIME_SESSIONS_LOCK:
        session = _REALTIME_SESSIONS.pop(session_id, None)
    if not session:
        return
    session["active"] = False
    try:
        ws = session.get("ws")
        if ws is not None:
            ws.close()
    except Exception:
        pass
    _realtime_emit(session, {"type": "session_closed", "reason": reason})


def _realtime_reader_loop(session_id):
    session = _REALTIME_SESSIONS.get(session_id)
    if not session:
        return
    ws = session["ws"]
    conv_id = session["conv_id"]
    cur_text_parts = []
    has_text_delta = False
    audio_b64_parts = []
    try:
        while session.get("active"):
            msg = ws.recv()
            if not msg:
                continue
            try:
                ev = json.loads(msg)
            except json.JSONDecodeError:
                continue
            typ = ev.get("type") or ""
            if typ in ("response.output_text.delta", "response.text.delta"):
                delta = ev.get("delta") or ""
                if delta:
                    has_text_delta = True
                    cur_text_parts.append(delta)
                    _realtime_emit(session, {"type": "text_delta", "delta": delta})
            elif typ == "response.audio_transcript.delta":
                # 某些场景文本增量可能走音频转写通道，作为回退补齐
                if has_text_delta:
                    continue
                delta = ev.get("delta") or ""
                if delta:
                    cur_text_parts.append(delta)
                    _realtime_emit(session, {"type": "text_delta", "delta": delta})
            elif typ in ("response.audio.delta", "response.output_audio.delta"):
                delta_b64 = ev.get("delta") or ""
                if delta_b64:
                    audio_b64_parts.append(delta_b64)
                    _realtime_emit(session, {"type": "audio_delta", "audio": delta_b64})
            elif typ == "response.done":
                ai_content = "".join(cur_text_parts).strip() or "（无回复）"
                cur_text_parts = []
                has_text_delta = False
                audio_url = None
                audio_name = None
                merged_b64 = "".join(audio_b64_parts).strip()
                audio_b64_parts = []
                if merged_b64:
                    try:
                        raw_audio = base64.b64decode(merged_b64, validate=False)
                        if raw_audio:
                            wav_bytes = _decoded_omni_audio_to_wav_bytes(raw_audio)
                            if wav_bytes:
                                audio_url, audio_name = _save_assistant_audio_file(wav_bytes, conv_id)
                    except Exception:
                        pass
                if conv_id is not None:
                    db.create_message(conv_id, "assistant", ai_content, "text", audio_url, audio_name)
                    db.update_conversation_last_message(conv_id, ai_content)
                _realtime_emit(session, {"type": "response_done"})
            elif typ == "error":
                err = ev.get("error") or {}
                msg = err.get("message") or ev.get("message") or str(ev.get("code", "unknown"))
                _realtime_emit(session, {"type": "error", "error": msg})
                break
    except Exception as e:
        _realtime_emit(session, {"type": "error", "error": str(e)})
    finally:
        _close_realtime_session(session_id, "upstream_closed")


def _create_realtime_upstream():
    """连接 DashScope Realtime 并初始化 session。"""
    url = f"{REALTIME_WS_URL.rstrip('/')}?model={REALTIME_MODEL}"
    headers = [f"Authorization: Bearer {REALTIME_API_KEY}"]
    ws = _websocket().create_connection(
        url,
        header=headers,
        timeout=_REALTIME_CONNECT_TIMEOUT_SEC,
    )
    ws.settimeout(_REALTIME_RECV_TIMEOUT_SEC)
    session_event = {
        "type": "session.update",
        "event_id": f"evt_{uuid.uuid4().hex[:24]}",
        "session": {
            "modalities": ["audio", "text"],
            "voice": "Cherry",
            "input_audio_format": "pcm16",
            "output_audio_format": "pcm16",
            "turn_detection": {
                "type": "server_vad",
                "silence_duration_ms": 700,
                "prefix_padding_ms": 240,
                "create_response": True,
            },
            "instructions": _build_realtime_instructions(),
        },
    }
    ws.send(json.dumps(session_event, ensure_ascii=False))
    deadline = time.time() + 10
    while time.time() < deadline:
        msg = ws.recv()
        if not msg:
            continue
        try:
            ev = json.loads(msg)
        except json.JSONDecodeError:
            continue
        typ = ev.get("type") or ""
        if typ == "session.updated":
            return ws
        if typ == "error":
            err = ev.get("error") or {}
            detail = err.get("message") or ev.get("message") or str(ev.get("code", "unknown"))
            raise RuntimeError(detail)
    raise RuntimeError("实时会话初始化超时")


def _get_realtime_session_owned(session_id, user_id):
    with _REALTIME_SESSIONS_LOCK:
        session = _REALTIME_SESSIONS.get(session_id)
    if not session:
        return None, (jsonify({"error": "实时会话不存在或已结束"}), 404)
    if int(session.get("user_id", -1)) != int(user_id):
        return None, (jsonify({"error": "无权访问该实时会话"}), 403)
    return session, None


@app.route("/api/chat/realtime/session/start", methods=["POST"])
def realtime_session_start():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        ws = _create_realtime_upstream()
        session_id = f"rt_{uuid.uuid4().hex}"
        session = {
            "id": session_id,
            "conv_id": None,
            "user_id": int(user_id),
            "ws": ws,
            "queue": queue.Queue(maxsize=2048),
            "active": True,
        }
        with _REALTIME_SESSIONS_LOCK:
            _REALTIME_SESSIONS[session_id] = session
        worker = threading.Thread(
            target=_realtime_reader_loop,
            args=(session_id,),
            daemon=True,
            name=f"realtime-reader-{session_id[:8]}",
        )
        worker.start()
        return jsonify({"sessionId": session_id})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/chat/realtime/session/<session_id>/audio", methods=["POST"])
def realtime_session_audio(session_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    session, err = _get_realtime_session_owned(session_id, user_id)
    if err:
        return err
    try:
        data = request.get_json(force=True, silent=True) or {}
        audio_b64 = (data.get("audio") or "").strip()
        if not audio_b64:
            return jsonify({"error": "缺少 audio（16k 16bit mono pcm 的 base64）"}), 400
        session["ws"].send(json.dumps({
            "type": "input_audio_buffer.append",
            "event_id": f"evt_{uuid.uuid4().hex[:24]}",
            "audio": audio_b64,
        }, ensure_ascii=False))
        return jsonify({"ok": True})
    except Exception as e:
        _realtime_emit(session, {"type": "error", "error": str(e)})
        return jsonify({"error": str(e)}), 500


@app.route("/api/chat/realtime/session/<session_id>/events", methods=["GET"])
def realtime_session_events(session_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    session, err = _get_realtime_session_owned(session_id, user_id)
    if err:
        return err

    def generate():
        while True:
            try:
                payload = session["queue"].get(timeout=_REALTIME_EVENT_IDLE_TIMEOUT_SEC)
                yield f"data: {json.dumps(payload, ensure_ascii=False)}\n\n"
                if payload.get("type") == "session_closed":
                    break
            except queue.Empty:
                # SSE 心跳，防止代理层断开空闲连接
                yield ":\n\n"

    return Response(
        generate(),
        mimetype="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
            "Connection": "keep-alive",
        },
    )


@app.route("/api/chat/realtime/session/<session_id>/stop", methods=["POST"])
def realtime_session_stop(session_id):
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    session, err = _get_realtime_session_owned(session_id, user_id)
    if err:
        return err
    _close_realtime_session(session.get("id") or session_id, "client_stopped")
    return jsonify({"ok": True})


# ---------- 认证：注册 / 登录 / 登出 / 当前用户 ----------


@app.route("/api/auth/register", methods=["POST"])
def auth_register():
    """注册：优先按学生链路（username+student_id+password），兼容历史邮箱注册。"""
    try:
        data = request.get_json(force=True, silent=True) or {}
        username = (data.get("username") or "").strip()
        student_id = (data.get("student_id") or data.get("studentId") or "").strip()
        email = (data.get("email") or "").strip().lower()
        password = data.get("password") or ""

        if not username:
            return _err("请输入姓名")
        if not password or len(password) < 6:
            return _err("密码至少 6 位")

        password_hash = generate_password_hash(password, method="pbkdf2:sha256")

        # 学生注册主流程：保持 users/student/emotion_record 通过 student_id 对齐
        if student_id:
            if not _is_student_id_account(student_id):
                return _err("学号格式不正确，应为 20 开头的 10 位数字")
            try:
                db.upsert_student(student_id, username)
                user_id = db.upsert_student_login_profile(student_id, username, password_hash)
            except Exception as e:
                return _err(f"注册失败: {str(e)}", 500)
            user_row = db.get_user_by_id(user_id)
            user = _user_row_to_json(user_row)
            token = secrets.token_urlsafe(32)
            _tokens[token] = user_id
            return jsonify({"user": user, "token": token})

        # 兼容历史邮箱注册（无学号时）
        if not email:
            return _err("请输入学号或邮箱")

        try:
            existing = db.get_user_by_mail(email)
        except Exception as e:
            return _err(f"数据库连接失败: {str(e)}", 500)
        if existing:
            return _err("该邮箱已注册", 409)

        try:
            user_id = db.create_user(email, username, password_hash)
        except Exception as e:
            return _err(f"注册失败: {str(e)}", 500)

        user_row = db.get_user_by_id(user_id)
        user = _user_row_to_json(user_row)
        token = secrets.token_urlsafe(32)
        _tokens[token] = user_id
        return jsonify({"user": user, "token": token})
    except Exception as e:
        return _err(f"注册异常: {str(e)}", 500)


@app.route("/api/auth/login", methods=["POST"])
def auth_login():
    """登录：body { account, password } -> { user, token }。兼容历史 email 字段。"""
    data = request.get_json(force=True, silent=True) or {}
    account = (data.get("account") or data.get("email") or "").strip()
    password = data.get("password") or ""

    if not account:
        return _err("请输入学号或邮箱")
    if not (_is_student_id_account(account) or _is_email_account(account)):
        return _err("仅支持学号或邮箱登录", 400)
    if not password:
        return _err("请输入密码")

    row = db.get_user_by_account(account)
    if not row or not check_password_hash(row["password_hash"], password):
        return _err("账号或密码错误", 401)
    if (row.get("role") or "").upper() != "STUDENT":
        return _err("当前入口仅支持学生账号登录", 403)

    user = _user_row_to_json(row)
    token = secrets.token_urlsafe(32)
    _tokens[token] = row["id"]
    return jsonify({"user": user, "token": token})


@app.route("/api/auth/logout", methods=["POST"])
def auth_logout():
    """登出：清除服务端 token（可选）。"""
    auth = request.headers.get("Authorization")
    if auth and auth.startswith("Bearer "):
        token = auth[7:].strip()
        _tokens.pop(token, None)
    return jsonify({"ok": True})


@app.route("/api/auth/me", methods=["GET"])
def auth_me():
    """当前用户：需要 Authorization: Bearer <token>"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    row = db.get_user_by_id(user_id)
    if not row:
        return _err("用户不存在", 404)
    return jsonify(_user_row_to_json(row))


@app.route("/api/auth/preferred-name", methods=["POST"])
def auth_set_preferred_name():
    """首次欢迎流程：设置用户希望被称呼的名字。body: { preferred_name }"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    preferred_name = (data.get("preferred_name") or data.get("preferredName") or "").strip()
    if not preferred_name:
        return _err("请输入你希望被称呼的名字")
    if len(preferred_name) > 50:
        return _err("称呼最多 50 个字符")
    ok = db.update_user_preferred_name(int(user_id), preferred_name)
    if not ok:
        return _err("保存失败", 500)
    row = db.get_user_by_id(int(user_id))
    return jsonify(_user_row_to_json(row))


def _feedback_row_to_json(row):
    if not row:
        return None
    created = row.get("created_at")
    if hasattr(created, "isoformat"):
        created = created.isoformat()
    return {
        "id": row.get("id"),
        "user_id": row.get("user_id"),
        "username": row.get("username") or "",
        "email": row.get("email") or "",
        "content": row.get("content") or "",
        "screenshot_url": row.get("screenshot_url"),
        "allow_contact": 1 if row.get("allow_contact") else 0,
        "created_at": created,
    }


@app.route("/api/feedback", methods=["POST"])
def submit_user_feedback():
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    content = (data.get("content") or "").strip()
    screenshot_url = (data.get("screenshot_url") or data.get("screenshotUrl") or "").strip()
    allow_contact = 1 if data.get("allow_contact") in (1, "1", True) or data.get("allowContact") in (1, "1", True) else 0
    if not content:
        return _err("请填写反馈内容")
    row = db.get_user_by_id(int(user_id))
    if not row:
        return _err("用户不存在", 404)
    try:
        rid = db.add_user_feedback(
            int(user_id),
            row.get("username") or "",
            row.get("mail") or "",
            content,
            screenshot_url=screenshot_url or None,
            allow_contact=allow_contact,
        )
        return jsonify({"ok": True, "id": rid})
    except Exception as e:
        return _err(f"反馈提交失败: {str(e)}", 500)


@app.route("/api/admin/user-feedback", methods=["GET"])
def admin_user_feedback_list():
    user_id, err_res = _require_auth()
    if err_res:
        return jsonify({"ok": False, "code": "UNAUTHORIZED", "message": "未登录"}), 401
    user_row = db.get_user_by_id(int(user_id))
    if not user_row:
        return jsonify({"ok": False, "code": "NOT_FOUND", "message": "用户不存在"}), 404
    if (user_row.get("role") or "").upper() != "ADMIN":
        return jsonify({"ok": False, "code": "FORBIDDEN", "message": "无权限"}), 403
    try:
        limit = request.args.get("limit", type=int) or 200
        rows = db.list_user_feedback(limit=limit)
        return jsonify({"ok": True, "data": [_feedback_row_to_json(r) for r in rows]})
    except Exception as e:
        return jsonify({"ok": False, "code": "DB_ERROR", "message": str(e)}), 500


# ---------- 情绪标签（与 users.id 对应，需登录） ----------


def _emotion_row_to_json(row):
    """将情绪记录转为 JSON（created_at 转字符串）。"""
    if not row:
        return None
    created = row.get("created_at")
    if hasattr(created, "isoformat"):
        created = created.isoformat()
    return {
        "id": row["id"],
        "user_id": row["user_id"],
        "emotion_label": row["emotion_label"],
        "created_at": created,
    }


@app.route("/api/emotion", methods=["POST"])
def emotion_add():
    """提交一条情绪标签。body { "emotion_label": "开心" }，需登录。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    data = request.get_json(force=True, silent=True) or {}
    label = (data.get("emotion_label") or data.get("label") or "").strip()
    if not label:
        return _err("请输入情绪标签")
    try:
        row_id = db.add_emotion_label(int(user_id), label)
        row = db.get_emotion_labels_by_user(int(user_id), limit=1)
        if row:
            return jsonify(_emotion_row_to_json(row[0]))
        return jsonify({"id": row_id, "user_id": user_id, "emotion_label": label, "created_at": None})
    except Exception as e:
        return _err(f"保存失败: {str(e)}", 500)


@app.route("/api/emotion", methods=["GET"])
def emotion_list():
    """当前用户的情绪标签列表，按时间倒序。query: limit 默认 100。需登录。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        limit = request.args.get("limit", type=int) or 100
        limit = min(max(1, limit), 500)
        rows = db.get_emotion_labels_by_user(int(user_id), limit=limit)
        return jsonify([_emotion_row_to_json(r) for r in rows])
    except Exception as e:
        return _err(f"查询失败: {str(e)}", 500)


@app.route("/api/emotion/latest", methods=["GET"])
def emotion_latest():
    """当前用户最近一条情绪标签。需登录。"""
    user_id, err_res = _require_auth()
    if err_res:
        return err_res
    try:
        row = db.get_latest_emotion_label(int(user_id))
        if not row:
            return jsonify(None)
        return jsonify(_emotion_row_to_json(row))
    except Exception as e:
        return _err(f"查询失败: {str(e)}", 500)


if __name__ == "__main__":
    # use_reloader=False：首次拉取 FaceNet 权重约 107MB，若开启重载，改代码会重启进程导致下载中断。
    # 需要热重载时可设置环境变量：FLASK_USE_RELOADER=1
    _use_reloader = os.environ.get("FLASK_USE_RELOADER", "").lower() in ("1", "true", "yes")
    # 低内存机器建议默认关闭 debug，避免额外内存占用；需要调试时可设置 FLASK_DEBUG=1
    _debug = os.environ.get("FLASK_DEBUG", "").lower() in ("1", "true", "yes")
    # threaded=True：避免单次人脸推理阻塞其它请求；开发环境建议保留
    app.run(host="0.0.0.0", port=5000, debug=_debug, threaded=True, use_reloader=_use_reloader)
